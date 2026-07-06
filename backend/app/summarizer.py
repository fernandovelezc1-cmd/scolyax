from __future__ import annotations

import re
import string
import math
from collections import Counter, defaultdict
from typing import Tuple

# Lazy imports - only loaded when needed
try:
    from PyPDF2 import PdfReader
except ImportError:
    PdfReader = None

try:
    from docx import Document as DocxDocument
except ImportError:
    DocxDocument = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

STOPWORDS = {
    "a","al","algo","algunas","algunos","ante","antes","como","con","contra","cual","cuales","cuando","de",
    "del","desde","donde","dos","el","ella","ellas","ellos","en","entre","era","erais","eran","eras","eres",
    "es","esa","esas","ese","eso","esos","esta","estaba","estado","estais","estamos","estan","estar","este",
    "estos","estoy","fin","fue","fueron","fui","fuimos","ha","hace","hacen","hacer","hacia","han","hasta",
    "incluso","la","las","le","les","lo","los","mas","me","mi","mis","mucho","muy","nada","ni","no","nos",
    "nosotros","o","otra","otras","otro","otros","para","pero","poco","por","porque","que","se","sea","segun",
    "ser","si","sin","sobre","sois","solamente","solo","somos","son","soy","su","sus","tambien","teneis",
    "tenemos","tengo","ti","tiene","tienen","todo","tras","tu","tus","un","una","uno","y","ya",
    "an","the","and","or","but","if","while","of","in","on","for","to","from","by","with","as","is",
    "are","was","were","be","been","being","at","it","its","this","that","these","those","i","you","he",
    "she","we","they","them","me","my","your","their","our","not","no","so","do","does","did"
}

PUNCT = set(string.punctuation + "¿¡""«»…")

# Palabras que indican oraciones importantes
IMPORTANCE_MARKERS = {
    "importante", "fundamental", "esencial", "clave", "principal", "critico", "vital",
    "necesario", "crucial", "significativo", "destacado", "relevante", "basico",
    "primordial", "indispensable", "conclusion", "resumen", "finalmente", "por tanto",
    "en consecuencia", "por lo tanto", "en conclusion", "important", "key", "essential",
    "critical", "fundamental", "vital", "necessary", "crucial", "significant", "main",
    "primary", "conclusion", "summary", "finally", "therefore", "consequently"
}

def _sentences(text: str) -> list[str]:
    """Divide el texto en oraciones con mejor precisión y manejo de archivos grandes."""
    text = re.sub(r"\s+", " ", text).strip()
    
    # Limita el texto procesado a 1 millón de caracteres para archivos muy grandes
    # Esto mantiene la precisión mientras maneja archivos enormes
    if len(text) > 1000000:
        text = text[:1000000]
    
    # Mejorado: maneja abreviaturas comunes
    text = re.sub(r"\b([Dd]r|[Dd]ra|[Ss]r|[Ss]ra|[Pp]rof|[Pp]hd|etc|p\.ej)\.", r"\1<PERIOD>", text)
    
    # Divide por puntuación final seguida de espacios
    parts = re.split(r"(?<=[\.\?\!])\s+|—\s+", text)
    
    # Restaura los puntos de abreviaturas y limpia
    parts = [s.replace("<PERIOD>", ".").strip() for s in parts if s.strip()]
    
    # Filtra oraciones demasiado cortas o que claramente no son oraciones
    filtered_parts = []
    for part in parts:
        # Mantiene oraciones con al menos 3 palabras o que claramente son válidas
        word_count = len(part.split())
        if word_count >= 3 or (word_count >= 2 and part[-1] in ".!?"):
            filtered_parts.append(part)
    
    return filtered_parts

def _tokens(text: str) -> list[str]:
    """Tokeniza el texto preservando contexto."""
    text = text.lower()
    text = "".join(ch for ch in text if ch not in PUNCT)
    return [t for t in re.split(r"\W+", text) if t and t not in STOPWORDS and not t.isdigit()]

def _calculate_tf_idf(sents: list[str]) -> dict[int, dict[str, float]]:
    """Calcula TF-IDF para cada oración."""
    # Term Frequency por oración
    sent_tf: list[dict[str, int]] = []
    for s in sents:
        toks = _tokens(s)
        sent_tf.append(Counter(toks))
    
    # Document Frequency
    df: dict[str, int] = defaultdict(int)
    for tf_dict in sent_tf:
        for term in tf_dict.keys():
            df[term] += 1
    
    # TF-IDF por oración
    num_sents = len(sents)
    tf_idf: dict[int, dict[str, float]] = {}
    
    for i, tf_dict in enumerate(sent_tf):
        tf_idf[i] = {}
        max_freq = max(tf_dict.values()) if tf_dict else 1
        
        for term, freq in tf_dict.items():
            # TF normalizado
            tf = freq / max_freq
            # IDF
            idf = math.log(num_sents / (df[term] + 1)) + 1
            tf_idf[i][term] = tf * idf
    
    return tf_idf

def _calculate_sentence_similarity(sent_a_tf_idf: dict[str, float], sent_b_tf_idf: dict[str, float]) -> float:
    """Calcula similitud del coseno entre dos oraciones."""
    common_terms = set(sent_a_tf_idf.keys()) & set(sent_b_tf_idf.keys())
    if not common_terms:
        return 0.0
    
    numerator = sum(sent_a_tf_idf[term] * sent_b_tf_idf[term] for term in common_terms)
    
    sum_a = sum(v**2 for v in sent_a_tf_idf.values())
    sum_b = sum(v**2 for v in sent_b_tf_idf.values())
    
    denominator = math.sqrt(sum_a) * math.sqrt(sum_b)
    
    return numerator / denominator if denominator > 0 else 0.0

def _score_sentences(sents: list[str]) -> list[tuple[int, float]]:
    """
    Algoritmo mejorado de scoring que combina:
    - TF-IDF para relevancia de términos
    - Similitud con otras oraciones (cohesión)
    - Posición en el documento
    - Longitud óptima
    - Marcadores de importancia
    - Número de palabras clave
    - Densidad de información
    """
    if not sents:
        return []
    
    # Calcula TF-IDF
    tf_idf = _calculate_tf_idf(sents)
    
    # Calcula scores base TF-IDF
    base_scores = []
    for i, sent in enumerate(sents):
        if i in tf_idf:
            score = sum(tf_idf[i].values())
        else:
            score = 0.0
        base_scores.append(score)
    
    # Normaliza scores base
    max_base = max(base_scores) if base_scores else 1.0
    if max_base > 0:
        base_scores = [s / max_base for s in base_scores]
    
    # Calcula matriz de similitud para cohesión
    similarity_scores = []
    for i in range(len(sents)):
        avg_similarity = 0.0
        if i in tf_idf:
            similarities = []
            for j in range(len(sents)):
                if i != j and j in tf_idf:
                    sim = _calculate_sentence_similarity(tf_idf[i], tf_idf[j])
                    similarities.append(sim)
            avg_similarity = sum(similarities) / len(similarities) if similarities else 0.0
        similarity_scores.append(avg_similarity)
    
    # Calcula score de posición (inicio y final son más importantes)
    position_scores = []
    num_sents = len(sents)
    for i in range(num_sents):
        if i == 0:  # Primera oración (introduce el tema)
            pos_score = 1.0
        elif i < 2:  # Primeras oraciones (contexto)
            pos_score = 0.85
        elif i >= num_sents - 1:  # Última oración (conclusión)
            pos_score = 0.8
        elif i >= num_sents - 3:  # Penúltimas oraciones (conclusión)
            pos_score = 0.7
        elif i <= 3:  # Primeras 4 oraciones después del intro
            pos_score = 0.6
        else:  # Oraciones del medio
            pos_score = 0.4
        position_scores.append(pos_score)
    
    # Analiza longitud, marcadores y densidad de información
    scored = []
    for i, sent in enumerate(sents):
        toks = _tokens(sent)
        num_toks = len(toks)
        
        # Bonus por longitud óptima (entre 10 y 25 palabras es muy bueno)
        if 10 <= num_toks <= 25:
            length_bonus = 0.4
        elif 8 <= num_toks <= 30:
            length_bonus = 0.2
        elif 6 <= num_toks <= 35:
            length_bonus = 0.1
        else:
            length_bonus = 0.0
        
        # Bonus por marcadores de importancia
        importance_bonus = 0.0
        sent_lower = sent.lower()
        marker_count = sum(1 for marker in IMPORTANCE_MARKERS if marker in sent_lower)
        if marker_count > 0:
            importance_bonus = min(0.35, 0.15 * marker_count)
        
        # Penalización por ser demasiado corta o muy larga
        if num_toks < 4:
            length_bonus = -0.2
        elif num_toks > 40:
            length_bonus = max(-0.1, length_bonus - 0.1)
        
        # Score de densidad: cuantas palabras únicas / total
        if num_toks > 0:
            unique_words = len(set(toks))
            density_score = (unique_words / num_toks) * 0.3
        else:
            density_score = 0.0
        
        # Score final combinado con pesos optimizados
        final_score = (
            base_scores[i] * 0.4 +           # TF-IDF relevancia (aumentado)
            similarity_scores[i] * 0.2 +      # Cohesión con otras oraciones
            position_scores[i] * 0.15 +       # Posición en el documento
            length_bonus * 0.15 +             # Longitud óptima
            importance_bonus +                # Marcadores de importancia
            density_score                     # Densidad de información
        )
        
        scored.append((i, final_score))
    
    return scored

async def generate_summary(file, text: str | None, sentences: int = 5) -> tuple[str, str]:
    extracted = ""
    if file is not None:
        fname = (file.filename or "").lower()
        try:
            if fname.endswith(".pdf"):
                if PdfReader is None:
                    raise ImportError("PyPDF2 is required for PDF processing. Install it with: pip install PyPDF2")
                reader = PdfReader(file.file)
                pages = []
                # Limita a 500 páginas para archivos muy grandes
                for idx, page in enumerate(reader.pages):
                    if idx >= 500:
                        break
                    text_page = page.extract_text() or ""
                    if text_page.strip():
                        pages.append(text_page)
                extracted = "\n".join(pages)
            elif fname.endswith(".docx"):
                if DocxDocument is None:
                    raise ImportError("python-docx is required for DOCX processing. Install it with: pip install python-docx")
                doc = DocxDocument(file.file)
                extracted = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
            elif fname.endswith(".pptx"):
                if Presentation is None:
                    raise ImportError("python-pptx is required for PPTX processing. Install it with: pip install python-pptx")
                pres = Presentation(file.file)
                buf = []
                for slide in pres.slides:
                    slide_text = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_text.append(shape.text)
                    if slide_text:
                        buf.append(" ".join(slide_text))
                extracted = "\n".join(buf)
            elif fname.endswith(".md") or fname.endswith(".markdown"):
                # Markdown: extrae solo el contenido de texto, sin formato
                raw = file.file.read() or b""
                try:
                    text_content = raw.decode("utf-8", errors="ignore")
                except Exception:
                    text_content = raw.decode(errors="ignore")
                # Limpia markdown: elimina caracteres especiales pero mantiene estructura
                import re
                text_content = re.sub(r'[#*_`\[\](){}<>!\\|-]', '', text_content)
                extracted = text_content
            elif fname.endswith(".rtf"):
                # RTF: extrae texto plano sin formato
                raw = file.file.read() or b""
                try:
                    text_content = raw.decode("utf-8", errors="ignore")
                except Exception:
                    text_content = raw.decode(errors="ignore")
                # Limpia RTF: elimina códigos RTF
                import re
                text_content = re.sub(r'\\[a-z]+\d*\s*', '', text_content)
                text_content = re.sub(r'[{}]', '', text_content)
                extracted = text_content
            else:
                # Asume que es archivo de texto plano
                raw = file.file.read() or b""
                try:
                    extracted = raw.decode("utf-8", errors="ignore")
                except Exception:
                    extracted = raw.decode(errors="ignore")
        except Exception as e:
            raise ValueError(f"Error procesando archivo '{fname}': {str(e)}")

    base = extracted or (text or "")
    summary, original = await summarize_from_text(base, sentences)
    return summary, original

async def summarize_from_text(text: str, sentences: int = 5) -> tuple[str, str]:
    original = text or ""
    if not original.strip():
        return "", ""
    
    # Limita a 500K caracteres para procesamiento (pero guarda el original completo)
    processing_text = original if len(original) <= 500000 else original[:500000]
    
    sents = _sentences(processing_text)
    
    # Si tiene muy pocas oraciones, retorna todo
    if len(sents) <= sentences:
        return " ".join(sents), original
    
    # Para textos muy largos, aumenta el número de oraciones a procesar
    # para asegurar mejor cobertura del documento
    if len(sents) > 200:
        adjusted_sentences = min(sentences * 2, len(sents) // 10)
    else:
        adjusted_sentences = sentences
    
    scored = _score_sentences(sents)
    
    # Selecciona las mejores oraciones manteniendo el orden original
    top_idx = {i for i, _ in sorted(scored, key=lambda x: x[1], reverse=True)[:adjusted_sentences]}
    
    # Mantiene el orden original para mejor legibilidad
    ordered = [s for i, s in enumerate(sents) if i in top_idx]
    
    return " ".join(ordered), original

def top_keywords(text: str, k: int = 10) -> list[str]:
    """
    Extrae palabras clave usando TF-IDF mejorado con análisis de n-gramas y filtrado inteligente.
    Prioriza frases significativas sobre palabras individuales.
    """
    toks = _tokens(text or "")
    if not toks:
        return []
    
    # Frecuencias de unigrams
    unigram_freq = Counter(toks)
    
    # Bigramas (pares de palabras consecutivas)
    bigrams = [" ".join(pair) for pair in zip(toks, toks[1:])]
    bigram_freq = Counter(bigrams)
    
    # Trigramas (tres palabras consecutivas)
    trigrams = [" ".join(triple) for triple in zip(toks, toks[1:], toks[2:])]
    trigram_freq = Counter(trigrams)
    
    # Cuatrigramas para frases clave (hasta 4 palabras)
    four_grams = [" ".join(quad) for quad in zip(toks, toks[1:], toks[2:], toks[3:])]
    four_gram_freq = Counter(four_grams)
    
    results = []
    
    # Procesa cuatrigramas (máxima prioridad)
    for quad, freq in four_gram_freq.most_common(k * 2):
        if freq >= 2:  # Al menos 2 apariciones
            score = freq * 4.0
            results.append((quad, score))
    
    # Procesa trigramas (alta prioridad)
    for trigram, freq in trigram_freq.most_common(k * 2):
        if freq >= 2:
            score = freq * 3.0
            results.append((trigram, score))
    
    # Procesa bigramas (media prioridad)
    for bigram, freq in bigram_freq.most_common(k * 3):
        if freq >= 2:
            score = freq * 2.0
            results.append((bigram, score))
    
    # Procesa unigrams (baja prioridad, pero más específicos)
    for word, freq in unigram_freq.most_common(k * 4):
        # Bonus para palabras más largas (tienden a ser más específicas)
        if len(word) > 4:
            length_bonus = (len(word) - 4) / 10.0
        else:
            length_bonus = 0.0
        score = freq * 1.0 + length_bonus
        results.append((word, score))
    
    # Ordena por score y elimina duplicados inteligentemente
    seen = set()
    unique_results = []
    
    for keyword, score in sorted(results, key=lambda x: x[1], reverse=True):
        # Evita keywords que sean substrings significativos de otros
        is_subset = False
        keyword_words = set(keyword.split())
        
        for existing in seen:
            existing_words = set(existing.split())
            # Si es substring exacto, salta
            if keyword in existing or existing in keyword:
                is_subset = True
                break
            # Si más del 70% de las palabras coinciden, probablemente es muy similar
            if len(keyword_words) > 0 and len(existing_words) > 0:
                overlap = len(keyword_words & existing_words)
                overlap_ratio = overlap / min(len(keyword_words), len(existing_words))
                if overlap_ratio > 0.7:
                    is_subset = True
                    break
        
        if not is_subset:
            unique_results.append(keyword)
            seen.add(keyword)
        
        if len(unique_results) >= k:
            break
    
    # Si no hay suficientes keywords, rellena con los más frecuentes
    if len(unique_results) < k:
        for w, _ in unigram_freq.most_common(k * 2):
            if w not in unique_results and w not in seen:
                unique_results.append(w)
                if len(unique_results) >= k:
                    break
    
    return unique_results[:k]

def extract_keywords(text: str, limit: int = 10) -> list[str]:
    """Alias para top_keywords para compatibilidad con código existente."""
    return top_keywords(text, limit)

def summarize_text(text: str, sentences: int = 5) -> list[str]:
    """
    Versión síncrona para compatibilidad.
    Retorna lista de oraciones del resumen.
    """
    if not text.strip():
        return []
    sents = _sentences(text)
    if len(sents) <= sentences:
        return sents
    scored = _score_sentences(sents)
    top_idx = {i for i, _ in sorted(scored, key=lambda x: x[1], reverse=True)[:sentences]}
    ordered = [s for i, s in enumerate(sents) if i in top_idx]
    return ordered
