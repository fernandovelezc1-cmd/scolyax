# ════════════════════════════════════════════════════════════════════
# CONFIGURACIÓN INICIAL: carga las variables del archivo .env antes
# de importar cualquier otro módulo, para que claves como SUPABASE_URL
# y los Client IDs de OAuth estén disponibles globalmente.
# ════════════════════════════════════════════════════════════════════
from pathlib import Path
from dotenv import load_dotenv
load_dotenv(dotenv_path=Path(__file__).resolve().parents[1] / '.env')

"""Punto de entrada principal de la API de Scolyax.

Este módulo inicializa la aplicación de FastAPI, configura el middleware de CORS y
define todos los endpoints necesarios para operar la plataforma estudiantil
Scolyax. Las rutas abarcan registro e inicio de sesión, tareas, recordatorios,
horarios semanales, sesiones de enfoque, estadísticas del tablero y el servicio
de resúmenes.

Cada función incluye docstrings descriptivos en español para facilitar el
mantenimiento por parte del equipo y asegurar que cualquier persona comprenda el
objetivo de cada pieza de lógica sin necesidad de explorar otros archivos.
"""

# ────────────────────────────────────────────────────────────────────
# LIBRERÍA ESTÁNDAR: utilidades de logging, concurrencia, HTTP, sistema
# de archivos, identificadores únicos, fechas y manipulación de URLs.
# ────────────────────────────────────────────────────────────────────
import logging
import asyncio
import httpx
import os
import uuid
from datetime import datetime, timezone
from typing import Optional, List
from urllib.parse import parse_qsl, urlencode, urljoin, urlparse, urlunparse

# ────────────────────────────────────────────────────────────────────
# FRAMEWORK WEB: FastAPI para definir endpoints y middleware;
# Pydantic para validar y serializar datos de entrada/salida.
# ────────────────────────────────────────────────────────────────────
from fastapi import BackgroundTasks, Body, FastAPI, File, Form, HTTPException, Query, Request, Response, UploadFile, Header
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import RedirectResponse, StreamingResponse
from pydantic import BaseModel

# ────────────────────────────────────────────────────────────────────
# MÓDULOS INTERNOS DE SCOLYAX:
#   summarizer          → resúmenes de texto con NLP/IA
#   calendar_integration→ crea/elimina eventos en Google y Microsoft
#   notebooklm_*        → integración con Google Gemini para sugerencias
#   tts_provider        → síntesis de voz (Google Cloud TTS / gTTS)
#   google_calendar     → cliente directo de la API de Google Calendar
#   ai_time_management  → herramientas IA para gestión del tiempo
#   mailer              → envío de correos + sistema de outbox offline
#   admin               → router con endpoints de administración
# ────────────────────────────────────────────────────────────────────
from . import summarizer
from . import calendar_integration
from . import notebooklm_integration
from . import notebooklm_tasks
from . import tts_provider
from . import google_calendar
from . import ai_time_management
from .mailer import send_registration_email, send_reminder_email, send_signin_email, resend_outbox
from .admin import router as admin_router
# ────────────────────────────────────────────────────────────────────
# MODELOS PYDANTIC: estructuras de datos que representan cada entidad
# del dominio (tarea, recordatorio, sesión, usuario, estadísticas...).
# Pydantic valida automáticamente tipos y formatos al recibir JSON.
# ────────────────────────────────────────────────────────────────────
from .models import (
    AIStudySession,
    AITimeManagementTool,
    AIToolRecommendation,
    AuthProvider,
    Checkpoint,
    CheckpointVerificationRequest,
    CheckpointVerificationResponse,
    CrisisSession,
    CrisisSessionCreate,
    DashboardStats,
    DisplayNameUpdate,
    EnergyEntry,
    EnergyEntryCreate,
    EnergyLevel,
    FocusSession,
    MicroTask,
    Reminder,
    ReminderCreate,
    ReminderUpdate,
    ScheduleEntry,
    Session,
    SessionCreate,
    SummaryRequest,
    SummaryResponse,
    Task,
    TaskStatus,
    TestAnalysisResult,
    User,
    OnboardingComplete,
    UserStatsPayload,
    UserFeedback,
    PomodoroEstimateRequest,
    PomodoroEstimateResponse,
    TimeEstimateResponse,
)
# ────────────────────────────────────────────────────────────────────
# CAPA DE ALMACENAMIENTO: todas las operaciones de lectura/escritura
# en Supabase pasan por este módulo. NUNCA se hacen queries directas
# fuera de aquí. El email del usuario es la clave primaria de aislamiento.
# ────────────────────────────────────────────────────────────────────
from .storage import (
    compute_dashboard_stats,
    load_focus_sessions,
    load_focus_sessions_for_user,
    load_reminders,
    load_users,
    load_sessions,
    load_schedule,
    load_tasks,
    next_id,
    save_focus_sessions,
    save_reminders,
    save_sessions,
    save_schedule,
    save_stats,
    save_tasks,
    save_users,
    save_token_for_email,
    load_tokens,
    create_user_session,
    validate_session_token,
    get_session_by_email,
    invalidate_session,
    invalidate_all_sessions,
    update_session_onboarding,
    get_supabase,
)

# ────────────────────────────────────────────────────────────────────
# FEEDBACK DE USUARIO: guarda valoraciones de logros y funcionalidades
# para que el equipo pueda analizar qué herramientas son más útiles.
# ────────────────────────────────────────────────────────────────────
from .user_feedback import (
    save_user_feedback,
    check_user_feedback_exists,
    load_all_user_feedback,
    get_feedback_stats,
)

# ────────────────────────────────────────────────────────────────────
# MÓDULO OAUTH: gestiona el flujo de autenticación con Google y
# Microsoft. OAuthStateStore almacena tokens de estado temporales
# (TTL 60 s) para prevenir ataques CSRF. is_stub_mode() permite
# pruebas locales sin credenciales reales.
# ────────────────────────────────────────────────────────────────────
from .oauth import (
    OAuthStateStore,
    get_oauth_client,
    get_google_client,
    get_google_calendar_client,
    is_stub_mode,
    resolve_frontend_base_url,
)

from contextlib import asynccontextmanager

# ════════════════════════════════════════════════════════════════════
# CICLO DE VIDA (LIFESPAN): se ejecuta al arrancar y apagar el servidor.
# Al arrancar:  corrige NULLs en notificaciones, inicia el scheduler
#               de emails de reactivación.
# Al apagar:    detiene el scheduler limpiamente.
# ════════════════════════════════════════════════════════════════════
async def lifespan(app: FastAPI):
    """Gestiona el ciclo de vida de la aplicación (startup/shutdown)"""
    # Startup
    logger.info("🚀 Iniciando aplicación Scolyax...")

    # Migración: actualizar notificaciones con sent=NULL a sent=false
    # (bug corregido: las notificaciones se insertaban sin sent explícito,
    #  y el query .eq("sent", False) no capturaba NULLs en PostgreSQL)
    try:
        from .supabase_client import get_supabase_client
        _sb = get_supabase_client()
        _sb.table("scheduled_notifications").update(
            {"sent": False}
        ).is_("sent", "null").execute()
        logger.info("✅ Migración: notificaciones con sent=NULL corregidas a sent=false")
    except Exception as e:
        logger.debug(f"Migración sent=NULL: {e}")
    
    # Iniciar scheduler para emails de reactivación
    try:
        from .scheduler import start_scheduler
        start_scheduler()
        logger.info("✅ Scheduler de emails de reactivación iniciado")
    except Exception as e:
        logger.error(f"❌ Error iniciando scheduler: {e}")
    
    yield
    
    # Shutdown
    logger.info("🛑 Deteniendo aplicación Scolyax...")
    try:
        from .scheduler import stop_scheduler
        stop_scheduler()
        logger.info("✅ Scheduler detenido correctamente")
    except Exception as e:
        logger.error(f"❌ Error deteniendo scheduler: {e}")

# ════════════════════════════════════════════════════════════════════
# INICIALIZACIÓN DE LA APP FASTAPI:
#   - lifespan controla startup/shutdown
#   - se incluye el router del panel de administración
#   - logger global para registrar eventos en todos los endpoints
# ════════════════════════════════════════════════════════════════════
app = FastAPI(title="Scolyax API", version="1.0.0", lifespan=lifespan)

# Incluir routers
app.include_router(admin_router)

logger = logging.getLogger(__name__)

# Contador en memoria de resúmenes generados (se reinicia con el servidor)
_summaries_generated: int = 0


# ════════════════════════════════════════════════════════════════════
# PROCESADOR DE NOTIFICACIONES EN BACKGROUND:
# _background_process_notifications() se llama desde endpoints
# frecuentes (GET /session, GET /tasks) para procesar push pendientes
# sin necesidad de un cron job externo. El cooldown de 30 s evita
# saturar Supabase con llamadas repetidas.
# ════════════════════════════════════════════════════════════════════
import threading
import time as _time

_last_bg_process_time = 0.0
_BG_PROCESS_COOLDOWN = 30  # segundos entre procesados para no saturar


def _background_process_notifications():
    """Procesa notificaciones push pendientes en un hilo secundario.
    Se invoca desde endpoints frecuentes para asegurar que los push
    se envíen incluso si el scheduler del dyno estaba dormido."""
    global _last_bg_process_time
    now = _time.time()
    if now - _last_bg_process_time < _BG_PROCESS_COOLDOWN:
        return  # cooldown activo, no re-procesar
    _last_bg_process_time = now
    try:
        from .notification_scheduler import process_pending_notifications
        sent = process_pending_notifications()
        if sent > 0:
            logger.info(f"📬 Background: {sent} notificaciones push enviadas al despertar")
    except Exception as e:
        logger.debug(f"Background push process: {e}")

# ════════════════════════════════════════════════════════════════════
# MIDDLEWARE HTTP:
#   1. normalize_double_slashes → corrige URLs del tipo //ai/generate
#      que algunos proxies o clientes generan incorrectamente.
#   2. CORSMiddleware → permite que el frontend (Vercel) pueda llamar
#      a la API desde un dominio distinto. En producción se podría
#      restringir allow_origins a solo el dominio de Vercel.
# ════════════════════════════════════════════════════════════════════
# Middleware para normalizar barras dobles en URLs (//ai/generate → /ai/generate)
@app.middleware("http")
async def normalize_double_slashes(request, call_next):
    import re
    if "//" in request.scope["path"]:
        request.scope["path"] = re.sub(r"/+", "/", request.scope["path"])
    return await call_next(request)


# Configuración de CORS para producción y desarrollo
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["GET", "POST", "PUT", "DELETE", "PATCH", "OPTIONS"],
    allow_headers=["*"],
    max_age=3600,
)


# ════════════════════════════════════════════════════════════════════
# CACHÉ EN MEMORIA:
# Para reducir la cantidad de queries a Supabase se implementan tres
# cachés con TTL configurable:
#   _users_cache      (TTL 5 min) → lista de usuarios registrados
#   _reminders_cache  (TTL 30 s)  → recordatorios por email
#   _tokens_cache     (TTL 60 s)  → tokens OAuth almacenados
# Cada caché tiene funciones _get_cached_*() e _invalidate_*_cache()
# que se llaman en los endpoints correspondientes.
# ════════════════════════════════════════════════════════════════════
ALLOWED_DOMAINS: dict[AuthProvider, tuple[str, ...]] = {
    # Empty tuples allow ANY email domain for that provider
    # Set to empty tuple to accept all users
    AuthProvider.GOOGLE: (),  # Allow all Gmail and any Google account
    AuthProvider.MICROSOFT: (),  # Allow all Outlook/Microsoft accounts
}

oauth_state_store = OAuthStateStore()

# Simple in-memory cache for users to avoid repeated Supabase queries during login
_users_cache = None
_users_cache_time = None
_CACHE_TTL_SECONDS = 300  # 5 minutos (mejorado de 60s para logins simultáneos)


def _get_cached_users():
    """Get users from cache or load from Supabase if cache expired."""
    global _users_cache, _users_cache_time
    
    import time
    now = time.time()
    
    # If cache exists and not expired, return it
    if _users_cache is not None and _users_cache_time is not None:
        if now - _users_cache_time < _CACHE_TTL_SECONDS:
            logger.debug(f"✅ Using cached users (age: {now - _users_cache_time:.1f}s)")
            return _users_cache
    
    # Cache expired or empty, load from Supabase
    logger.info("📥 Loading users from Supabase (cache miss or expired)")
    _users_cache = load_users()
    _users_cache_time = now
    return _users_cache


def _invalidate_users_cache():
    """Invalidate the users cache."""
    global _users_cache, _users_cache_time
    _users_cache = None
    _users_cache_time = None
    logger.debug("🗑️ Users cache invalidated")


def save_users_cached(users):
    """Save users to storage and invalidate cache."""
    from .storage import save_users as _save_users
    _save_users(users)
    _invalidate_users_cache()
    logger.info(f"💾 Saved {len(users)} users and invalidated cache")


# Simple in-memory cache for reminders and tokens (similar to users cache)
_reminders_cache = {}  # {email: (reminders_list, timestamp)}
_tokens_cache = None
_tokens_cache_time = None
_CACHE_TTL_REMINDERS = 30  # Cache reminders for 30 seconds
_CACHE_TTL_TOKENS = 60  # Cache tokens for 60 seconds


def _get_cached_reminders(email: str):
    """Get reminders from cache or load from Supabase if cache expired."""
    global _reminders_cache
    
    import time
    now = time.time()
    
    # If cache exists and not expired, return it
    if email in _reminders_cache:
        reminders_list, cache_time = _reminders_cache[email]
        if now - cache_time < _CACHE_TTL_REMINDERS:
            logger.debug(f"✅ Using cached reminders for {email} (age: {now - cache_time:.1f}s)")
            return reminders_list
    
    # Cache expired or empty, load from Supabase
    logger.info(f"📥 Loading reminders from Supabase for {email} (cache miss or expired)")
    from .storage import load_reminders as _load_reminders
    reminders_list = _load_reminders(email)
    _reminders_cache[email] = (reminders_list, now)
    return reminders_list


def _invalidate_reminders_cache(email: str = None):
    """Invalidate reminders cache for a specific email or all."""
    global _reminders_cache
    if email:
        if email in _reminders_cache:
            del _reminders_cache[email]
            logger.debug(f"🗑️ Reminders cache invalidated for {email}")
    else:
        _reminders_cache = {}
        logger.debug("🗑️ All reminders cache invalidated")


def _get_cached_tokens():
    """Get tokens from cache or load from Supabase if cache expired."""
    global _tokens_cache, _tokens_cache_time
    
    import time
    now = time.time()
    
    # If cache exists and not expired, return it
    if _tokens_cache is not None and _tokens_cache_time is not None:
        if now - _tokens_cache_time < _CACHE_TTL_TOKENS:
            logger.debug(f"✅ Using cached tokens (age: {now - _tokens_cache_time:.1f}s)")
            return _tokens_cache
    
    # Cache expired or empty, load from Supabase
    logger.info("📥 Loading tokens from Supabase (cache miss or expired)")
    from .storage import load_tokens as _load_tokens
    _tokens_cache = _load_tokens()
    _tokens_cache_time = now
    return _tokens_cache


def _invalidate_tokens_cache():
    """Invalidate the tokens cache."""
    global _tokens_cache, _tokens_cache_time
    _tokens_cache = None
    _tokens_cache_time = None
    logger.debug("🗑️ Tokens cache invalidated")


# ════════════════════════════════════════════════════════════════════
# FUNCIONES AUXILIARES GLOBALES:
#   _ensure_utc()            → normaliza fechas a UTC
#   _model_copy()            → compatibilidad Pydantic v1/v2 para clonar modelos
#   _resolve_redirect_target()→ valida URLs de retorno tras OAuth (anti-open-redirect)
#   _append_query_params()   → agrega parámetros a una URL preservando los existentes
#   _get_active_session()    → lee la sesión activa del disco (legado)
#   _require_session()       → valida el token Bearer; lanza 401 si no es válido
#   _normalize_email/name()  → limpieza de strings antes de guardar
#   _validate_email_provider()→ verifica que el email pertenezca al proveedor correcto
#   _find_user()             → busca usuario en caché por email normalizado
# ════════════════════════════════════════════════════════════════════
def _ensure_utc(dt: datetime) -> datetime:
    """Normaliza una marca de tiempo a UTC preservando la hora local."""
    if dt.tzinfo is None:
        # Si no tiene zona horaria, asumimos que es hora local
        local_tz = datetime.now().astimezone().tzinfo
        local_dt = dt.replace(tzinfo=local_tz)
        return local_dt.astimezone(timezone.utc)
    return dt.astimezone(timezone.utc)


def _model_copy(instance, **kwargs):
    """Obtiene una copia del modelo compatible con pydantic v1 y v2."""
    if hasattr(instance, "model_copy"):
        return instance.model_copy(**kwargs)  # type: ignore[call-arg]
    if hasattr(instance, "copy"):
        return instance.copy(**kwargs)  # type: ignore[call-arg]
    update = kwargs.get("update") or {}
    payload = {
        key: getattr(instance, key)
        for key in getattr(instance, "__dict__", {})
    }
    payload.update(update)
    return type(instance)(**payload)


def _resolve_redirect_target(next_url: Optional[str]) -> str:
    """Valida y construye la URL de retorno tras completar OAuth."""

    base = resolve_frontend_base_url(next_url)
    if not next_url:
        return base

    parsed_target = urlparse(next_url)
    parsed_base = urlparse(base)

    if parsed_target.scheme and parsed_target.netloc:
        if (
            parsed_target.netloc != parsed_base.netloc
            or parsed_target.scheme != parsed_base.scheme
        ):
            raise HTTPException(
                status_code=400,
                detail="La URL de retorno no coincide con el dominio configurado para el frontend.",
            )
        return next_url

    combined = urljoin(base + "/", next_url.lstrip("/"))
    return combined


def _append_query_params(url: str, params: dict[str, str]) -> str:
    """Agrega parámetros de consulta preservando los existentes."""

    parsed = urlparse(url)
    current = dict(parse_qsl(parsed.query, keep_blank_values=True))
    current.update(params)
    new_query = urlencode(current)
    return urlunparse(parsed._replace(query=new_query))


def _get_active_session() -> Session | None:
    """Obtiene la sesión activa almacenada en disco, si existe."""
    sessions = load_sessions()
    return sessions[0] if sessions else None


def _require_session(authorization: Optional[str] = None) -> Session:
    """Valida que haya una sesión activa mediante token o fallback a sesión en memoria.
    
    Args:
        authorization: Header Authorization con formato "Bearer <token>"
        
    Returns:
        Session objeto si la validación es exitosa
        
    Raises:
        HTTPException 401 si no hay sesión válida
    """
    # Prioridad 1: Validar token del header (sistema nuevo - producción)
    if authorization and authorization.startswith("Bearer "):
        session_token = authorization[7:]
        session_data = validate_session_token(session_token)
        
        if session_data:
            # Token válido, construir Session desde session_data
            user_data = session_data.get("user_data")
            if user_data:
                provider_str = user_data.get("provider", AuthProvider.GOOGLE.value)
                provider = AuthProvider(provider_str) if not isinstance(provider_str, AuthProvider) else provider_str
                
                return Session(
                    id=user_data["id"],
                    email=user_data["email"],
                    provider=provider,
                    display_name=user_data["display_name"],
                    session_token=session_token,
                )
            
            # Fallback: buscar usuario por email
            user = _find_user(session_data.get("email"))
            if user:
                return Session(
                    id=user.id,
                    email=user.email,
                    provider=user.provider,
                    display_name=user.display_name,
                    session_token=session_token,
                )
    
    # Prioridad 2: Fallback a sesión en memoria (desarrollo local)
    session = _get_active_session()
    if session is not None:
        return session
    
    # No hay sesión válida
    raise HTTPException(
        status_code=401,
        detail="Inicia sesión con Google o Microsoft para programar recordatorios y sincronizar notificaciones.",
    )


def _normalize_email(email: str) -> str:
    """Normaliza el correo eliminando espacios y pasando a minúsculas."""
    return email.strip().lower()


def _normalize_display_name(display_name: str, fallback_email: str) -> str:
    """Limpia el nombre mostrado y genera uno amigable si no se proporcionó."""
    normalized = display_name.strip()
    if normalized:
        return " ".join(normalized.split())
    return fallback_email.split("@", 1)[0].replace(".", " ").title()


def _validate_email_provider(email: str, provider: AuthProvider) -> None:
    """Comprueba que el correo pertenezca al proveedor (Google/Microsoft) esperado."""
    email = email.strip()
    if "@" not in email:
        raise HTTPException(status_code=400, detail="El correo electrónico debe contener un dominio válido.")
    domain = email.split("@", 1)[1].lower()
    
    # Para Google, validar que sea Gmail SOLO si ALLOWED_DOMAINS no está vacío
    if provider == AuthProvider.GOOGLE:
        allowed = ALLOWED_DOMAINS.get(provider, ())
        # Si allowed está vacío, permitir todos los dominios
        if allowed and domain not in allowed:
            raise HTTPException(
                status_code=400,
                detail="Utiliza un correo de Gmail para conectarte con Google.",
            )
    # Para Microsoft, permitir cualquier dominio (cuentas personales, educativas, empresariales)
    # ya que Microsoft OAuth valida que sea una cuenta válida de Microsoft


def _find_user(email: str) -> User | None:
    """Busca un usuario persistido usando el correo normalizado.
    Usa caché en memoria para evitar múltiples queries a Supabase durante login.
    """
    normalized_email = _normalize_email(email)
    users = _get_cached_users()  # Use cache instead of load_users()
    for user in users:
        if user.email == normalized_email:
            return user
    return None


# ════════════════════════════════════════════════════════════════════
# HELPERS DE SESIÓN Y NOTIFICACIÓN:
#   _send_reminder_notification() → envía solo el email del recordatorio;
#                                    los push los gestiona el scheduler.
#   _persist_session_for_user()   → crea el token de sesión en Supabase
#                                    y lo devuelve junto con el objeto Session.
# ════════════════════════════════════════════════════════════════════
def _send_reminder_notification(reminder: Reminder) -> None:
    """Envía solo el correo del recordatorio, sin interrumpir la API.
    Las push notifications se gestionan via notification_scheduler."""

    session = _get_active_session()
    if session is None:
        return
    user = _find_user(session.email)
    if user is None:
        return
    try:
        send_reminder_email(reminder, user)
    except Exception as exc:  # noqa: BLE001 - queremos registrar fallos inesperados
        logger.warning("No se pudo enviar el correo del recordatorio: %s", exc)


def _persist_session_for_user(user: User, user_agent: Optional[str] = None, ip_address: Optional[str] = None) -> Session:
    """Crea y guarda la sesión activa en Supabase y devuelve la sesión con token."""
    try:
        # Create session in Supabase and get token
        session_token = create_user_session(user.email, user_agent=user_agent, ip_address=ip_address)
        logger.warning(f"[TOKEN] Session token created for {user.email}: {session_token[:20]}...")
    except Exception as e:
        logger.error(f"❌ Could not create Supabase session for {user.email}: {e}")
        session_token = None
    
    session = Session(
        id=user.id,
        email=user.email,
        provider=user.provider,
        display_name=user.display_name,
    )
    
    # Also save to in-memory sessions for backward compatibility
    save_sessions([session])
    
    # Add session_token to response (will be used by frontend)
    session.session_token = session_token
    logger.warning(f"Session object prepared: {session.email}, token: {session.session_token[:20] if session.session_token else 'None'}...")
    
    return session


# TEMPORARILY DISABLED: Background tasks causing startup issues
# TODO: Re-enable after fixing async/await handling in startup
# The reminder dispatcher and email resender have been disabled to resolve startup crashes
# These features will be re-enabled in a future update


# ════════════════════════════════════════════════════════════════════
# ENDPOINTS DE SALUD (HEALTH CHECKS):
#   GET /         → ping básico, confirma que el servidor responde
#   GET /health   → verifica que la API está corriendo
#   GET /health/ready → readiness probe de Railway (acepta tráfico)
#   GET /health/live  → liveness probe de Railway (servidor vivo)
# Railway usa estos endpoints para decidir si el contenedor está listo
# y si debe reiniciarlo automáticamente en caso de fallo.
# ════════════════════════════════════════════════════════════════════
@app.get("/")
def root() -> dict[str, str]:
    """Root endpoint for basic connectivity check."""
    return {"status": "ok", "api": "Scolyax"}


@app.get("/health")
def health_check() -> dict[str, str]:
    """Permite verificar que la API está viva y disponible."""
    return {"status": "ok", "message": "Scolyax API is running"}


@app.get("/health/ready")
def health_ready() -> dict[str, str]:
    """Railway readiness probe - returns 200 if service is ready to accept traffic."""
    return {"status": "ready", "message": "Scolyax API is ready"}


@app.get("/health/live")
def health_live() -> dict[str, str]:
    """Railway liveness probe - returns 200 if service is alive."""
    return {"status": "live", "message": "Scolyax API is alive"}


# ═══════════════════════════════════════════════════════
# ══ PUSH NOTIFICATIONS ════════════════════════════════
# ═══════════════════════════════════════════════════════

# ════════════════════════════════════════════════════════════════════
# PUSH NOTIFICATIONS (Web Push / VAPID):
#   GET  /push/vapid-public-key  → devuelve la clave pública VAPID para
#                                   que el SW del frontend pueda suscribirse
#   POST /push/subscribe         → guarda la suscripción del navegador
#   POST /push/unsubscribe       → elimina una suscripción
#   POST /push/test              → envía una push de prueba inmediata
#   GET  /push/debug             → diagnóstico detallado del par de claves VAPID
#   GET  /push/status            → resumen del estado del sistema push
#   GET  /push/scheduled         → lista notificaciones pendientes/enviadas
#   POST /push/process-now       → procesa pendientes de inmediato (testing)
#   GET  /push/scheduler-status  → estado del scheduler interno
# ════════════════════════════════════════════════════════════════════

@app.get("/push/vapid-public-key")
def get_vapid_public_key():
    """Retorna la clave pública VAPID para que el frontend se suscriba a push."""
    from .push_notifications import get_public_vapid_key, is_push_available
    key = get_public_vapid_key()
    return {
        "available": is_push_available(),
        "publicKey": key
    }


@app.post("/push/subscribe")
def subscribe_push(
    subscription: dict = Body(...),
    authorization: Optional[str] = Header(None)
):
    """Guarda la suscripción push del usuario actual."""
    from .push_notifications import save_push_subscription
    session = _require_session(authorization)
    logger.info(f"📬 Push subscribe: {session.email}, endpoint: {str(subscription.get('endpoint', ''))[:60]}")
    success = save_push_subscription(session.email, subscription)
    if success:
        return {"status": "ok", "message": "Suscripción push guardada"}
    logger.error(f"❌ Push subscribe failed for {session.email}")
    raise HTTPException(status_code=500, detail="Error guardando suscripción push")


@app.post("/push/unsubscribe")
def unsubscribe_push(
    data: dict = Body(...),
    authorization: Optional[str] = Header(None)
):
    """Elimina una suscripción push."""
    from .push_notifications import remove_push_subscription
    _require_session(authorization)  # Verificar sesión
    endpoint = data.get("endpoint", "")
    if not endpoint:
        raise HTTPException(status_code=400, detail="Falta el endpoint")
    remove_push_subscription(endpoint)
    return {"status": "ok", "message": "Suscripción push eliminada"}


@app.post("/push/test")
def test_push_notification(
    authorization: Optional[str] = Header(None)
):
    """Envía una notificación push de prueba al usuario. Devuelve diagnóstico por suscripción."""
    from .push_notifications import get_user_subscriptions, is_push_available, _send_push, _classify_endpoint
    session = _require_session(authorization)

    available = is_push_available()
    subs = get_user_subscriptions(session.email)
    logger.info(f"📨 Push test: user={session.email}, vapid={available}, subs={len(subs)}")

    if not available:
        return {"status": "error", "sent": 0, "total_subs": 0,
                "reason": "VAPID no configurado en el servidor"}
    if not subs:
        return {"status": "error", "sent": 0, "total_subs": 0,
                "reason": "No hay suscripciones push guardadas para este usuario. Presiona 'Resuscribir Push' primero."}

    results = []
    sent_count = 0
    for sub in subs:
        endpoint = sub.get("endpoint", "")
        service = _classify_endpoint(endpoint)
        ok = _send_push(
            sub,
            title="🧠 Scolyax – Prueba",
            body="¡Push recibido correctamente! Si ves esto con la app cerrada, todo funciona.",
            tag="push-test",
            url="/"
        )
        if ok:
            sent_count += 1
        results.append({
            "service": service,
            "endpoint_prefix": endpoint[:50] + "..." if len(endpoint) > 50 else endpoint,
            "ok": ok,
        })
        logger.info(f"📨 Push test → {service}: {'✅' if ok else '❌'} ({endpoint[:60]})")

    return {
        "status": "ok" if sent_count > 0 else "error",
        "sent": sent_count,
        "total_subs": len(subs),
        "results": results,
    }


@app.get("/push/debug")
def push_debug(authorization: Optional[str] = Header(None)):
    """Diagnóstico detallado: muestra qué hay en la suscripción almacenada y prueba VAPID."""
    import os
    from .push_notifications import get_user_subscriptions, _get_vapid_keys
    session = _require_session(authorization)
    subs = get_user_subscriptions(session.email)

    pub, priv, email = _get_vapid_keys()

    # Validar par de claves VAPID
    vapid_valid = False
    vapid_error = None
    derived_public_key = None
    keys_match = None
    try:
        from py_vapid import Vapid
        import base64
        from cryptography.hazmat.primitives.serialization import Encoding, PublicFormat
        import time as _time

        vv = Vapid.from_string(private_key=priv)
        # Derivar clave pública desde la privada
        pub_bytes = vv.public_key.public_bytes(Encoding.X962, PublicFormat.UncompressedPoint)
        derived_public_key = base64.urlsafe_b64encode(pub_bytes).rstrip(b'=').decode()

        # Comparar con la clave pública configurada
        keys_match = (derived_public_key == pub) if pub else None

        # Intentar firmar un claim de prueba
        test_claims = {
            "sub": email if email.startswith("mailto:") else f"mailto:{email}",
            "aud": "https://fcm.googleapis.com",
            "exp": int(_time.time()) + 3600,
        }
        headers = vv.sign(test_claims)
        vapid_valid = bool(headers.get("Authorization") or headers.get("authorization"))
    except Exception as ve:
        vapid_error = str(ve)

    vapid_info = {
        "public_key_len": len(pub) if pub else 0,
        "public_key_prefix": pub[:25] + "..." if pub else None,
        "private_key_len": len(priv) if priv else 0,
        "private_key_prefix": priv[:10] + "..." if priv else None,
        "claims_email": email,
        "has_mailto": email.startswith("mailto:") if email else False,
        "key_pair_valid": vapid_valid,
        "key_pair_error": vapid_error,
        "derived_public_key_prefix": derived_public_key[:30] + "..." if derived_public_key else None,
        "configured_public_key_prefix": pub[:30] + "..." if pub else None,
        "keys_match": keys_match,
    }

    sub_details = []
    for i, sub in enumerate(subs):
        endpoint = sub.get("endpoint", "")
        keys = sub.get("keys", {})
        sub_details.append({
            "index": i,
            "endpoint_prefix": endpoint[:80] + "..." if len(endpoint) > 80 else endpoint,
            "endpoint_service": "FCM" if "fcm.googleapis.com" in endpoint
                               else "Mozilla" if "push.services.mozilla.com" in endpoint
                               else "other",
            "has_keys": bool(keys),
            "has_p256dh": bool(keys.get("p256dh")),
            "p256dh_len": len(keys.get("p256dh", "")),
            "has_auth": bool(keys.get("auth")),
            "auth_len": len(keys.get("auth", "")),
            "extra_fields": [k for k in sub.keys() if k not in ("endpoint", "keys", "expirationTime")],
        })

    return {
        "user_email": session.email,
        "vapid": vapid_info,
        "subscriptions_count": len(subs),
        "subscriptions": sub_details,
    }


@app.get("/push/status")
def push_status(authorization: Optional[str] = Header(None)):
    """Diagnóstico completo del sistema de push notifications."""
    from .push_notifications import is_push_available, get_user_subscriptions
    from .supabase_client import get_supabase_client

    session = _require_session(authorization)
    status = {
        "vapid_configured": is_push_available(),
        "user_email": session.email,
        "subscriptions": 0,
        "scheduled_pending": 0,
        "scheduled_sent": 0,
        "tables_ok": True,
        "issues": [],
    }

    try:
        supabase = get_supabase_client()
    except Exception:
        status["tables_ok"] = False
        status["issues"].append("No se pudo conectar a Supabase")
        return status

    # Verificar tabla push_subscriptions
    try:
        subs = get_user_subscriptions(session.email)
        status["subscriptions"] = len(subs)
    except Exception:
        status["tables_ok"] = False
        status["issues"].append("Tabla push_subscriptions no existe — ejecuta el SQL de migración")

    # Verificar tabla scheduled_notifications
    try:
        pending = supabase.table("scheduled_notifications").select(
            "id", count="exact"
        ).eq("user_email", session.email).eq("sent", False).execute()
        status["scheduled_pending"] = pending.count if pending.count is not None else len(pending.data or [])

        sent_resp = supabase.table("scheduled_notifications").select(
            "id", count="exact"
        ).eq("user_email", session.email).eq("sent", True).execute()
        status["scheduled_sent"] = sent_resp.count if sent_resp.count is not None else len(sent_resp.data or [])
    except Exception:
        status["tables_ok"] = False
        status["issues"].append("Tabla scheduled_notifications no existe — ejecuta el SQL de migración")

    if not status["vapid_configured"]:
        status["issues"].append("Claves VAPID no configuradas en el servidor (VAPID_PUBLIC_KEY / VAPID_PRIVATE_KEY)")
    if status["subscriptions"] == 0:
        status["issues"].append("No hay suscripciones push para este usuario — activa notificaciones en la app")

    return status


@app.get("/push/scheduled")
def get_scheduled_notifications(authorization: Optional[str] = Header(None)):
    """Lista las notificaciones programadas del usuario actual (pendientes + recientes enviadas)."""
    from .supabase_client import get_supabase_client

    session = _require_session(authorization)

    try:
        supabase = get_supabase_client()
        # Pendientes
        pending = supabase.table("scheduled_notifications").select("*").eq(
            "user_email", session.email
        ).eq("sent", False).order("send_at", desc=False).limit(50).execute()

        # Últimas 20 enviadas
        sent = supabase.table("scheduled_notifications").select("*").eq(
            "user_email", session.email
        ).eq("sent", True).order("sent_at", desc=True).limit(20).execute()

        return {
            "pending": pending.data or [],
            "sent": sent.data or [],
        }
    except Exception as e:
        logger.warning("Error obteniendo notificaciones programadas: %s", e)
        return {"pending": [], "sent": [], "error": "Tabla scheduled_notifications no disponible"}


@app.post("/push/process-now")
def process_notifications_now(authorization: Optional[str] = Header(None)):
    """Fuerza el procesamiento inmediato de notificaciones pendientes (útil para testing)."""
    from .notification_scheduler import process_pending_notifications
    _require_session(authorization)
    sent = process_pending_notifications()
    return {"status": "ok", "sent": sent}


@app.get("/push/scheduler-status")
def get_scheduler_status_endpoint():
    """Diagnóstico: estado del scheduler de notificaciones."""
    try:
        from .scheduler import get_scheduler_status
        status = get_scheduler_status()
        return status
    except Exception as e:
        return {"status": "error", "error": str(e)}


# ════════════════════════════════════════════════════════════════════
# GESTIÓN DE SESIÓN:
#   GET  /session         → devuelve la sesión si el token Bearer es válido.
#                            Incluye datos de onboarding, herramientas
#                            recomendadas, racha y nivel de XP.
#   POST /session         → alias de /login para compatibilidad legacy.
#   POST /login           → inicio de sesión con email+proveedor.
#                            Si el usuario no existe, lo registra.
#   DELETE /session       → cierra sesión invalidando el token en Supabase.
#   PATCH /session/display-name → actualiza el nombre visible del usuario.
# ════════════════════════════════════════════════════════════════════
@app.get("/session", response_model=Optional[Session])
def get_session(authorization: Optional[str] = Header(None), background_tasks: BackgroundTasks = None) -> Session | None:
    """Devuelve la sesión actual si existe, de lo contrario `None`.
    
    Puede validar mediante:
    1. Token en header Authorization: "Bearer <session_token>"
    2. Sin token, devuelve None (no hay fallback a sesión en memoria)
    """
    # Procesar notificaciones push pendientes en background al abrir la app
    if background_tasks:
        background_tasks.add_task(_background_process_notifications)

    # Use debug-level for frequent health checks to avoid log spam in normal operation
    logger.debug(f"GET /session called with authorization: {authorization is not None}")

    # If no Authorization header, this is an anonymous request — return None (200)
    if not authorization:
        logger.debug("No Authorization header provided; returning None")
        return None

    # Authorization header present — expect 'Bearer <token>'
    if not authorization.startswith("Bearer "):
        logger.debug("Malformed Authorization header received")
        raise HTTPException(status_code=401, detail="Malformed Authorization header")

    session_token = authorization[7:]
    logger.debug(f"Validating session token: {session_token[:20]}...")
    session_data = validate_session_token(session_token)
    if not session_data:
        logger.debug("Token validation failed or token not found; returning 401")
        raise HTTPException(status_code=401, detail="Invalid or expired session token")

    # Check if user data is already included (optimized path)
    user_data = session_data.get("user_data")
    if user_data:
        # Fast path: user data already fetched during validation
        provider_str = user_data.get("provider", AuthProvider.GOOGLE.value)
        provider = AuthProvider(provider_str) if not isinstance(provider_str, AuthProvider) else provider_str
        
        # Parse recommended_tools if it's a string
        recommended_tools = user_data.get("recommended_tools")
        if isinstance(recommended_tools, str):
            try:
                import json
                recommended_tools = json.loads(recommended_tools)
            except:
                recommended_tools = []
        elif not isinstance(recommended_tools, list):
            recommended_tools = []
        
        session = Session(
            id=user_data["id"],
            email=user_data["email"],
            provider=provider,
            display_name=user_data["display_name"],
            session_token=session_token,
            has_completed_onboarding=user_data.get("has_completed_onboarding", False),
            selected_tool=user_data.get("selected_tool"),
            recommended_tools=recommended_tools,
            streak_days=user_data.get("streak_days", 0),
            total_xp=user_data.get("total_xp", 0),
            level=user_data.get("level", 1),
        )
        logger.info(f"Session restored (optimized) for {user_data['email']}")
        return session
    
    # Fallback: fetch user separately (slower path)
    user = _find_user(session_data.get("email"))
    if not user:
        logger.debug("Session token valid but user not found; returning 401")
        raise HTTPException(status_code=401, detail="User for token not found")

    # Parse recommended_tools if it's a string
    recommended_tools = getattr(user, "recommended_tools", None)
    if isinstance(recommended_tools, str):
        try:
            import json
            recommended_tools = json.loads(recommended_tools)
        except:
            recommended_tools = []
    elif not isinstance(recommended_tools, list):
        recommended_tools = []

    session = Session(
        id=user.id,
        email=user.email,
        provider=user.provider,
        display_name=user.display_name,
        session_token=session_token,
        has_completed_onboarding=getattr(user, "has_completed_onboarding", False),
        selected_tool=getattr(user, "selected_tool", None),
        recommended_tools=recommended_tools,
        streak_days=getattr(user, "streak_days", 0),
        total_xp=getattr(user, "total_xp", 0),
        level=getattr(user, "level", 1),
    )
    logger.info(f"Session restored from Supabase for {user.email}")
    return session


@app.post("/session", response_model=Session, status_code=201)
def create_session(payload: SessionCreate) -> Session:
    """Alias para `login` que permite compatibilidad con clientes antiguos."""
    return login(payload)


@app.post("/login", response_model=Session)
def login(payload: SessionCreate) -> Session:
    """Permite iniciar sesión; si el correo no existe, lo registra (auto-registro)."""
    normalized_email = _normalize_email(payload.email)
    user = _find_user(normalized_email)

    if user is None:
        _validate_email_provider(payload.email, payload.provider)
        from fastapi import Response
        display_name = _normalize_display_name(payload.display_name, normalized_email)
        response = Response()
        return _register_payload(
            SessionCreate(email=normalized_email, provider=payload.provider, display_name=display_name),
            response,
        )

    # Allow users to login with any provider (support multiple OAuth providers per email)
    # If they registered with Google, they can now also login with Microsoft and vice versa
    if user.provider != payload.provider:
        # Update the user's provider to the one being used for this login
        # This enables flexible multi-provider login
        users = _get_cached_users()  # Use cache
        for index, existing in enumerate(users):
            if existing.id == user.id:
                users[index] = _model_copy(user, update={"provider": payload.provider})
                user = users[index]
                save_users_cached(users)  # Invalidate cache
                break

    display_name = _normalize_display_name(payload.display_name, user.email)
    if display_name != user.display_name:
        user.display_name = display_name
        users = _get_cached_users()  # Use cache
        for index, existing in enumerate(users):
            if existing.id == user.id:
                users[index] = user
                break
        save_users_cached(users)  # Invalidate cache

    # Email de signin deshabilitado - Railway no puede conectar a SMTP
    # try:
    #     send_signin_email(user)
    # except Exception as exc:
    #     logger.warning('No se pudo enviar correo de inicio de sesión: %s', exc)
    return _persist_session_for_user(user)

@app.post("/auth/google/login", response_model=Session)
def login_google(payload: SessionCreate) -> Session:
    """Inicia sesión obligando a que el proveedor sea Google."""

    enforced = _enforce_provider(payload, AuthProvider.GOOGLE)
    return login(enforced)


@app.post("/auth/microsoft/login", response_model=Session)
def login_microsoft(payload: SessionCreate) -> Session:
    """Inicia sesión obligando a que el proveedor sea Microsoft."""

    enforced = _enforce_provider(payload, AuthProvider.MICROSOFT)
    return login(enforced)


def _register_payload(payload: SessionCreate, response: Response) -> Session:
    """Registra un nuevo usuario y envía el correo de bienvenida."""
    _validate_email_provider(payload.email, payload.provider)
    normalized_email = _normalize_email(payload.email)
    display_name = _normalize_display_name(payload.display_name, normalized_email)

    existing = _find_user(normalized_email)
    if existing is not None:
        # Allow users to switch providers - update their provider to the new one
        # This enables flexible multi-provider support without blocking users
        if existing.provider != payload.provider:
            users = _get_cached_users()  # Use cache
            for index, user in enumerate(users):
                if user.id == existing.id:
                    users[index] = _model_copy(user, update={
                        "provider": payload.provider,
                        "display_name": display_name
                    })
                    existing = users[index]
                    save_users_cached(users)  # Invalidate cache
                    break
        elif existing.display_name != display_name:
            users = _get_cached_users()  # Use cache
            for index, user in enumerate(users):
                if user.id == existing.id:
                    users[index] = _model_copy(user, update={"display_name": display_name})
                    existing = users[index]
                    break
            save_users_cached(users)  # Invalidate cache
        response.status_code = 200
        return _persist_session_for_user(existing)

    users = _get_cached_users()  # Use cache instead of load_users()
    user = User(
        id=next_id(users),
        email=normalized_email,
        provider=payload.provider,
        display_name=display_name,
        created_at=_ensure_utc(datetime.now(timezone.utc)),
    )
    users.append(user)
    save_users_cached(users)  # Invalidate cache

    response.status_code = 201
    # Emails deshabilitados temporalmente - Railway no puede conectar a SMTP
    # Cada intento tarda ~3 segundos en timeout, bloqueando el login
    # TODO: Implementar envío async o usar servicio de email cloud
    # try:
    #     send_registration_email(user)
    # except Exception as exc:
    #     logger.warning('No se pudo enviar correo de registro: %s', exc)
    # try:
    #     send_signin_email(user)
    # except Exception as exc:
    #     logger.warning('No se pudo enviar correo de inicio de sesión: %s', exc)
    return _persist_session_for_user(user)


def _enforce_provider(payload: SessionCreate, provider: AuthProvider) -> SessionCreate:
    """Crea un nuevo payload asegurando que el proveedor coincida con el endpoint."""

    return SessionCreate(email=payload.email, provider=provider, display_name=payload.display_name)


def _issue_oauth_state(
    provider: AuthProvider,
    mode: str,
    display_name: Optional[str],
    next_url: Optional[str],
    stub_email: Optional[str],
) -> str:
    """Genera y guarda el estado necesario para completar el flujo OAuth."""

    normalized_mode = "register" if mode == "register" else "login"
    payload = {
        "mode": normalized_mode,
        "provider": provider.value,
        "next": _resolve_redirect_target(next_url),
    }

    if display_name and display_name.strip():
        payload["display_name"] = " ".join(display_name.strip().split())

    if stub_email and stub_email.strip():
        payload["stub_email"] = _normalize_email(stub_email)

    return oauth_state_store.issue(payload)


# ════════════════════════════════════════════════════════════════════
# FLUJO OAUTH (autenticación federada):
#   GET /auth/google/start      → genera estado y redirige a Google OAuth
#   GET /auth/google/calendar/start → solicita scopes adicionales de Calendar
#   GET /auth/microsoft/start   → genera estado y redirige a Microsoft OAuth
#   GET /auth/google/callback   → recibe el código de Google, intercambia por
#                                  tokens, crea/actualiza usuario y sesion
#   GET /auth/microsoft/callback→ idéntico para Microsoft
#   POST /register              → registro directo (legacy, sin popup OAuth)
#   POST /auth/google/register  → registro forzando proveedor Google
#   POST /auth/microsoft/register→ registro forzando proveedor Microsoft
#
# Flujo resumido:
#   1. Frontend abre popup → llama a /auth/{provider}/start
#   2. Backend redirige al proveedor con `state` firmado (anti-CSRF)
#   3. Proveedor redirige a /auth/{provider}/callback con `code`
#   4. Backend intercambia code por tokens, obtiene perfil del usuario
#   5. Crea/actualiza usuario en Supabase y genera session_token
#   6. Redirige al frontend con ?auth=signed-in&token=...
# ════════════════════════════════════════════════════════════════════
@app.get("/auth/google/start")
def start_google_oauth(
    mode: str = "login",
    display_name: Optional[str] = None,
    next: Optional[str] = None,
    stub_email: Optional[str] = None,
) -> RedirectResponse:
    """Inicia el flujo OAuth de Google redirigiendo al consentimiento oficial."""

    state = _issue_oauth_state(AuthProvider.GOOGLE, mode, display_name, next, stub_email)
    client = get_oauth_client(AuthProvider.GOOGLE)

    if is_stub_mode():
        fake_callback = _append_query_params(
            client.config.redirect_uri,
            {"code": "stub-code", "state": state},
        )
        return RedirectResponse(fake_callback, status_code=307)

    prompt = "consent" if mode == "register" else None
    authorize_url = client.build_authorize_url(state, prompt=prompt)
    return RedirectResponse(authorize_url, status_code=307)


@app.get("/auth/google/calendar/start")
def start_google_calendar_oauth(
    session_token: Optional[str] = None,
    authorization: Optional[str] = Header(None),
    next: Optional[str] = None,
) -> RedirectResponse:
    """Inicia el flujo OAuth de Google con scopes de Calendar.
    
    El usuario ya debe estar autenticado. Acepta token vía query param o header.
    Esto solicita permisos adicionales de Google Calendar sin la pantalla de app no verificada.
    """
    # Aceptar token como query param (para redirects del browser) o como header
    if session_token:
        authorization = f"Bearer {session_token}"
    
    session = _require_session(authorization)
    
    # Crear estado con modo especial calendar_connect
    payload = {
        "mode": "calendar_connect",
        "provider": AuthProvider.GOOGLE.value,
        "next": _resolve_redirect_target(next),
        "email": session.email,
    }
    state = oauth_state_store.issue(payload)
    
    # Usar cliente de Google con scopes de Calendar
    client = get_google_calendar_client()
    
    if is_stub_mode():
        from .oauth import _build_config
        fake_callback = _append_query_params(
            client.config.redirect_uri,
            {"code": "stub-code", "state": state},
        )
        return RedirectResponse(fake_callback, status_code=307)
    
    authorize_url = client.build_authorize_url(state, prompt="consent")
    return RedirectResponse(authorize_url, status_code=307)


@app.get("/auth/microsoft/start")
def start_microsoft_oauth(
    mode: str = "login",
    display_name: Optional[str] = None,
    next: Optional[str] = None,
    stub_email: Optional[str] = None,
) -> RedirectResponse:
    """Inicia el flujo OAuth con Microsoft Azure (Outlook/Teams)."""

    state = _issue_oauth_state(AuthProvider.MICROSOFT, mode, display_name, next, stub_email)
    client = get_oauth_client(AuthProvider.MICROSOFT)

    if is_stub_mode():
        fake_callback = _append_query_params(
            client.config.redirect_uri,
            {"code": "stub-code", "state": state},
        )
        return RedirectResponse(fake_callback, status_code=307)

    authorize_url = client.build_authorize_url(state)
    return RedirectResponse(authorize_url, status_code=307)


@app.get("/auth/google/callback")
async def google_callback(
    code: Optional[str] = None,
    state: Optional[str] = None,
    error: Optional[str] = None,
) -> RedirectResponse:
    """Recibe la respuesta de Google y completa el flujo de autenticación."""

    if not state:
        raise HTTPException(status_code=400, detail="Falta el estado de autenticación de Google.")
    return await _complete_oauth_callback(AuthProvider.GOOGLE, code, state, error)


@app.get("/auth/microsoft/callback")
async def microsoft_callback(
    code: Optional[str] = None,
    state: Optional[str] = None,
    error: Optional[str] = None,
) -> RedirectResponse:
    """Recibe la respuesta de Microsoft y completa el flujo de autenticación."""

    if not state:
        raise HTTPException(status_code=400, detail="Falta el estado de autenticación de Microsoft.")
    return await _complete_oauth_callback(AuthProvider.MICROSOFT, code, state, error)


@app.post("/register", response_model=Session, status_code=201)
def register_user(payload: SessionCreate, response: Response) -> Session:
    """Endpoint genérico compatible con clientes antiguos."""

    return _register_payload(payload, response)


@app.post("/auth/google/register", response_model=Session, status_code=201)
def register_google(payload: SessionCreate, response: Response) -> Session:
    """Registra una cuenta garantizando que el proveedor sea Google."""

    enforced = _enforce_provider(payload, AuthProvider.GOOGLE)
    return _register_payload(enforced, response)


@app.post("/auth/microsoft/register", response_model=Session, status_code=201)
def register_microsoft(payload: SessionCreate, response: Response) -> Session:
    """Registra una cuenta garantizando que el proveedor sea Microsoft."""

    enforced = _enforce_provider(payload, AuthProvider.MICROSOFT)
    return _register_payload(enforced, response)


def _extract_email(profile: dict, provider: AuthProvider) -> str | None:
    """Obtiene el correo electrónico del perfil según el proveedor."""

    if provider is AuthProvider.GOOGLE:
        return profile.get("email")
    return (
        profile.get("mail")
        or profile.get("userPrincipalName")
        or profile.get("preferred_username")
        or profile.get("email")
    )


async def _complete_oauth_callback(
    provider: AuthProvider,
    code: Optional[str],
    state: str,
    error: Optional[str],
) -> RedirectResponse:
    """Procesa la respuesta del proveedor OAuth y redirige al frontend.
    
    Optimización: Incluye timing logs para medir performance de cada etapa.
    """
    import time
    start_time = time.time()
    logger.warning(f"[OAUTH] OAuth callback STARTED - Provider: {provider.value}")

    logger.info(f"OAuth callback received - Provider: {provider.value}, State: {state[:20]}..., Code present: {bool(code)}, Error: {error}")
    
    try:
        # Primero validamos sin consumir, por si Microsoft hace múltiples redirects
        state_start = time.time()
        state_payload = oauth_state_store.validate(state)
        logger.info(f"[OK] State validated successfully in {time.time() - state_start:.2f}s: {state_payload}")
        redirect_target = _resolve_redirect_target(state_payload.get("next"))
    except HTTPException as e:
        logger.error(f"State validation failed: {e.detail}")
        # Si la validación falla, redirigir al frontend con error
        redirect_target = resolve_frontend_base_url()
        params = {
            "auth": "error",
            "provider": provider.value,
            "message": "session_expired",
        }
        return RedirectResponse(_append_query_params(redirect_target, params), status_code=303)

    if error:
        # En caso de error, consumimos el state
        oauth_state_store.consume(state)
        params = {
            "auth": "error",
            "provider": provider.value,
            "message": error,
        }
        return RedirectResponse(_append_query_params(redirect_target, params), status_code=303)

    if not code:
        raise HTTPException(status_code=400, detail="El proveedor no entregó un código de autorización válido.")
    
    # Ahora que tenemos el código, consumimos el state para evitar reutilización
    oauth_state_store.consume(state)

    # Para calendar_connect, usamos el cliente con scopes de Calendar
    if state_payload.get("mode") == "calendar_connect":
        client = get_google_calendar_client()
    else:
        client = get_oauth_client(provider)
    try:
        logger.info(f"[EXCHANGE] Starting OAuth token exchange with {provider.value}...")
        # Exchange code for tokens
        exchange_start = time.time()
        tokens = await client.exchange_code(code)
        exchange_elapsed = time.time() - exchange_start
        logger.info(f"[OK] Code exchange completed in {exchange_elapsed:.2f}s")
        if not tokens or not tokens.get("access_token"):
            logger.error("No se recibieron tokens válidos del proveedor OAuth")
            raise HTTPException(status_code=400, detail="Error al obtener los tokens de autorización.")

        logger.info(f"[PROFILE] Starting profile fetch from {provider.value}...")
        # Fetch user profile
        profile_start = time.time()
        profile = await client.fetch_profile(tokens, state_payload)
        profile_elapsed = time.time() - profile_start
        logger.info(f"[OK] Profile fetch completed in {profile_elapsed:.2f}s")
        if not profile:
            logger.error("No se pudo obtener el perfil del usuario")
            raise HTTPException(status_code=400, detail="Error al obtener el perfil de usuario.")

        email = _extract_email(profile, provider)
        if not email:
            logger.error("El perfil no contiene un correo electrónico")
            raise HTTPException(status_code=400, detail="No se pudo obtener el correo electrónico del perfil autenticado.")
    except Exception as e:
        logger.exception("Error en el proceso de OAuth: %s", str(e))
        raise HTTPException(status_code=500, detail="Error durante el proceso de autenticación. Por favor, inténtelo de nuevo.")

    raw_display = state_payload.get("display_name") or (
        profile.get("name")
        or profile.get("displayName")
        or profile.get("given_name")
        or profile.get("preferred_username")
        or ""
    )
    display_name = _normalize_display_name(str(raw_display or ""), email)

    payload = SessionCreate(email=email, provider=provider, display_name=display_name)
    mode = state_payload.get("mode")
    
    session_result = None

    if mode == "calendar_connect":
        # Modo especial: solo guardar tokens de Calendar, sin login/register
        logger.info(f"[CALENDAR_CONNECT] Saving calendar tokens for {email}...")
        try:
            if tokens:
                save_token_for_email(email, tokens)
                logger.info(f"✅ Calendar tokens saved for {email}")
        except Exception as exc:
            logger.warning("No se pudieron guardar los tokens de Calendar: %s", exc)
        
        # Redirigir al frontend indicando que Calendar fue conectado
        params = {
            "calendar": "connected",
            "provider": provider.value,
        }
        final_url = _append_query_params(redirect_target, params)
        logger.info(f"[REDIRECT] Calendar connect redirect to: {final_url[:100]}...")
        total_time = time.time() - start_time
        logger.info(f"[TOTAL] Calendar connect completed in {total_time:.2f}s")
        return RedirectResponse(final_url, status_code=303)

    elif mode == "register":
        # Login or register
        logger.info(f"[REGISTER] Starting register for {email}...")
        auth_start = time.time()
        response = Response()
        if provider is AuthProvider.GOOGLE:
            session_result = register_google(payload, response)
        else:
            session_result = register_microsoft(payload, response)
        auth_elapsed = time.time() - auth_start
        logger.info(f"[OK] Register completed in {auth_elapsed:.2f}s")
        # Guardar tokens para el email asociado si hay refresh_token
        try:
            if tokens and tokens.get("refresh_token"):
                save_token_for_email(email, tokens)
        except Exception as exc:  # noqa: BLE001
            logger.warning("No se pudieron guardar los tokens de OAuth: %s", exc)
        status_label = "registered" if response.status_code == 201 else "signed-in"
    else:
        logger.info(f"[LOGIN] Starting login for {email}...")
        auth_start = time.time()
        if provider is AuthProvider.GOOGLE:
            session_result = login_google(payload)
        else:
            session_result = login_microsoft(payload)
        auth_elapsed = time.time() - auth_start
        logger.info(f"[OK] Login completed in {auth_elapsed:.2f}s")
        # Guardar tokens para ambos proveedores
        try:
            if tokens and tokens.get("refresh_token"):
                save_token_for_email(email, tokens)
        except Exception as exc:  # noqa: BLE001
            logger.warning("No se pudieron guardar los tokens de OAuth: %s", exc)
        status_label = "signed-in"

    params = {
        "auth": status_label,
        "provider": provider.value,
    }
    
    # Incluir el session token en los parámetros si está disponible
    if session_result and session_result.session_token:
        logger.warning(f"[TOKEN] Including token in redirect: {session_result.session_token[:20]}...")
        params["token"] = session_result.session_token
    else:
        logger.error(f"[ERROR] No token available! session_result: {session_result}, token: {session_result.session_token if session_result else 'None'}")
    
    final_url = _append_query_params(redirect_target, params)
    logger.warning(f"[REDIRECT] Redirecting to: {final_url[:100]}...")
    
    # Log total OAuth time
    total_time = time.time() - start_time
    logger.info(f"[TOTAL] TOTAL OAuth process completed in {total_time:.2f}s (Target: <3s)")
    if total_time > 5:
        logger.warning(f"[SLOW] OAuth took {total_time:.2f}s - Supabase connection may be slow")
    
    return RedirectResponse(final_url, status_code=303)


@app.delete("/session", status_code=204)
def clear_session(authorization: Optional[str] = Header(None)) -> None:
    """Cierra la sesión activa eliminando el registro en disco y en Supabase.
    
    Puede invalidar mediante:
    1. Token en header Authorization: "Bearer <session_token>"
    2. Sesión en memoria (legacy)
    """
    logger.warning(f"DELETE /session called with authorization: {authorization is not None}")
    
    # Intenta invalidar sesión Supabase si hay token
    if authorization and authorization.startswith("Bearer "):
        session_token = authorization[7:]
        try:
            invalidate_session(session_token)
            logger.info(f"✅ Invalidated Supabase session")
        except Exception as e:
            logger.warning(f"❌ Could not invalidate Supabase session: {e}")
    
    # Limpia sesión en memoria
    logger.warning(f"Clearing sessions.json - saving empty list")
    save_sessions([])
    logger.warning(f"✅ Sessions cleared.")


@app.patch("/session/display-name", response_model=Session)
def update_display_name(payload: DisplayNameUpdate, authorization: Optional[str] = Header(None)) -> Session:
    """Actualiza el nombre para mostrar de la sesión y del perfil guardado."""

    session = _require_session(authorization)
    normalized_name = _normalize_display_name(payload.display_name, session.email)
    users = load_users()

    for index, user in enumerate(users):
        if user.id != session.id:
            continue

        updated_user = _model_copy(user, update={"display_name": normalized_name})
        users[index] = updated_user
        save_users(users)
        return _persist_session_for_user(updated_user)

    raise HTTPException(status_code=404, detail="No se encontró el usuario asociado a la sesión activa")


# ════════════════════════════════════════════════════════════════════
# ONBOARDING / TEST COGNITIVO:
#   POST /api/onboarding-complete       → guarda la herramienta seleccionada
#                                          y las recomendadas al finalizar el
#                                          selector de herramientas.
#   POST /api/onboarding-test-completed → marca que el usuario terminó el
#                                          test cognitivo y guarda recommended_tools.
#   POST /dev/quick-login               → login rápido sin OAuth para desarrollo
#                                          local (bloqueado en producción).
# ════════════════════════════════════════════════════════════════════
@app.post("/api/onboarding-complete")
def complete_onboarding(payload: OnboardingComplete, authorization: Optional[str] = Header(None)) -> dict:
    """Guarda los datos del test cognitivo y la herramienta seleccionada."""
    
    session = _require_session(authorization)
    
    try:
        logger.info(f"💾 Starting onboarding completion for {session.email}")
        logger.info(f"   Selected tool: {payload.selected_tool}")
        logger.info(f"   Recommended tools: {payload.recommended_tools}")
        
        # Guardar en la base de datos
        update_session_onboarding(
            email=session.email,
            selected_tool=payload.selected_tool,
            recommended_tools=payload.recommended_tools
        )
        
        logger.info(f"✅ Onboarding saved to DB for {session.email}")
        logger.info(f"✅ Cache will be invalidated, next /session call will fetch fresh data")
        
        return {"success": True, "message": "Onboarding completed successfully"}
    except Exception as e:
        logger.error(f"❌ Failed to save onboarding for {session.email}: {e}")
        import traceback
        logger.error(traceback.format_exc())
        raise HTTPException(status_code=500, detail=f"Failed to save onboarding data: {str(e)}")


class TestCompletedPayload(BaseModel):
    """Payload para cuando el usuario completa el test cognitivo."""
    recommended_tools: List[str] = []


@app.post("/api/onboarding-test-completed")
def mark_test_completed(payload: TestCompletedPayload, authorization: Optional[str] = Header(None)) -> dict:
    """Marca que el usuario completó el test cognitivo.
    
    Esto guarda los recommended_tools pero SIN seleccionar herramienta todavía.
    La herramienta se selecciona después en /api/onboarding-complete
    """
    
    global _users_cache, _users_cache_time
    
    session = _require_session(authorization)
    
    try:
        logger.info(f"✅ Test completed for {session.email}")
        logger.info(f"   Recommended tools: {payload.recommended_tools}")
        
        # Guardar SOLO los recommended_tools y marcar test como completado
        # NO guardar selected_tool (será null hasta que el usuario lo seleccione)
        from .storage import get_supabase, invalidate_token_cache
        supabase = get_supabase()
        
        # Guardar como lista directamente (columna TEXT[] o TEXT con JSON)
        # Intentar primero como lista; si falla, guardar como JSON string
        tools_to_save = payload.recommended_tools if isinstance(payload.recommended_tools, list) else []
        try:
            response = supabase.table("users").update({
                "has_completed_onboarding": True,
                "recommended_tools": tools_to_save
            }).eq("email", session.email).execute()
        except Exception as save_err:
            import json
            logger.warning(f"⚠️ Failed to save as list, trying JSON string: {save_err}")
            response = supabase.table("users").update({
                "has_completed_onboarding": True,
                "recommended_tools": json.dumps(tools_to_save)
            }).eq("email", session.email).execute()
        
        logger.info(f"✅ Test marked as completed for {session.email}")
        logger.info(f"   Supabase response: {response}")
        
        # Invalidate all caches - both in main.py and storage.py
        invalidate_token_cache()  # Clear token cache in storage.py
        _users_cache = None  # Clear users cache in main.py
        _users_cache_time = None
        
        logger.info(f"🔄 Cache invalidated - next /session call will fetch fresh data from DB")
        
        return {"success": True, "message": "Test completed successfully"}
    except Exception as e:
        logger.error(f"❌ Failed to mark test as completed for {session.email}: {e}")
        import traceback
        logger.error(traceback.format_exc())
        raise HTTPException(status_code=500, detail=f"Failed to mark test as completed: {str(e)}")


@app.post("/dev/quick-login", response_model=Session)
def dev_quick_login(email: str) -> Session:
    """Endpoint de desarrollo para login rápido sin OAuth.
    
    Solo disponible cuando ENVIRONMENT != 'production'.
    Permite iniciar sesión directamente con un email sin pasar por OAuth.
    """
    environment = os.getenv("ENVIRONMENT", "development")
    if environment == "production":
        raise HTTPException(
            status_code=403,
            detail="Quick login no está disponible en producción. Usa OAuth."
        )
    
    normalized_email = _normalize_email(email)
    user = _find_user(normalized_email)
    
    if not user:
        raise HTTPException(
            status_code=404,
            detail=f"Usuario {email} no encontrado. Regístrate primero con OAuth."
        )
    
    logger.info(f"🚀 DEV Quick login for {user.email}")
    return _persist_session_for_user(user)


# ════════════════════════════════════════════════════════════════════
# TAREAS (CRUD completo):
#   GET    /tasks                      → lista todas las tareas del usuario
#   POST   /tasks                      → crea tarea y programa notificaciones push
#   POST   /tasks/estimate-pomodoros   → IA estima cuántos pomodoros necesita la tarea
#   PUT    /tasks/{id}                 → reemplaza completamente una tarea
#   PATCH  /tasks/{id}/status          → solo actualiza el estado (pending/in_progress/done)
#   DELETE /tasks/{id}                 → elimina tarea y cancela sus notificaciones
#   GET    /tasks/{id}/estimate-time   → IA estima el tiempo en minutos para la tarea
#
# La estimación de pomodoros/tiempo usa heurísticas por palabras clave
# (ensayo, proyecto, quiz...), longitud de notas y materia académica.
# ════════════════════════════════════════════════════════════════════
@app.get("/tasks", response_model=list[Task])
def list_tasks(authorization: Optional[str] = Header(None), background_tasks: BackgroundTasks = None) -> list[Task]:
    """Devuelve todas las tareas guardadas."""
    session = _require_session(authorization)
    # Procesar notificaciones pendientes en background al listar tareas
    if background_tasks:
        background_tasks.add_task(_background_process_notifications)
    return load_tasks(session.email)


@app.post("/tasks", response_model=Task, status_code=201)
def create_task(
    task: Task,
    background_tasks: BackgroundTasks,
    authorization: Optional[str] = Header(None),
    x_timezone: Optional[str] = Header(None, alias="X-Timezone"),
) -> Task:
    """Crea una nueva tarea asignando un identificador incremental."""
    session = _require_session(authorization)
    tasks = load_tasks(session.email)
    task.id = next_id(tasks)
    task.user_email = session.email  # Asignar email del usuario a la tarea
    tasks.append(task)
    save_tasks(tasks, user_email=session.email)

    # Programar notificaciones push en background (no bloquea la respuesta al cliente)
    def _schedule_notifications():
        try:
            from .notification_scheduler import schedule_task_notifications
            schedule_task_notifications(
                task_id=task.id,
                title=task.title,
                course=task.course,
                due_date=task.due_date,
                user_email=session.email,
                user_timezone=x_timezone,
            )
        except Exception as e:
            logger.warning("⚠️ No se pudieron programar notificaciones de tarea: %s", e)

    background_tasks.add_task(_schedule_notifications)

    return task


@app.post("/tasks/estimate-pomodoros", response_model=PomodoroEstimateResponse)
def estimate_pomodoros(request: PomodoroEstimateRequest, authorization: Optional[str] = Header(None)) -> PomodoroEstimateResponse:
    """Estima cuántos pomodoros tomará completar una tarea usando IA."""
    _require_session(authorization)  # Validar sesión
    
    # Analizar el contenido de la tarea
    title_lower = request.title.lower()
    notes_lower = request.notes.lower()
    course_lower = request.course.lower()
    
    # Palabras clave que indican tareas largas
    long_keywords = ['ensayo', 'proyecto', 'investigación', 'tesis', 'monografía', 'informe extenso', 
                     'presentación completa', 'análisis profundo', 'desarrollo completo']
    medium_keywords = ['tarea', 'ejercicios', 'práctica', 'leer capítulo', 'estudio', 'repaso',
                       'laboratorio', 'informe', 'resumen', 'presentación']
    short_keywords = ['quiz', 'test corto', 'lectura breve', 'repasar', 'revisar', 'corregir']
    
    # Calcular longitud del contenido
    notes_length = len(request.notes)
    
    # Análisis de urgencia
    days_until_due = None
    if request.due_date:
        try:
            days_until_due = (request.due_date - datetime.now(timezone.utc)).days
        except:
            days_until_due = None
    
    # Calcular estimación base
    estimated = 2  # Por defecto
    confidence = 0.7
    reasoning = ""
    suggestions = []
    
    # Ajuste por palabras clave
    if any(keyword in title_lower or keyword in notes_lower for keyword in long_keywords):
        estimated = 5
        confidence = 0.85
        reasoning = f"📚 Tarea de tipo académico extenso detectado. Requiere múltiples sesiones de concentración profunda."
        suggestions.append("Divide la tarea en secciones (introducción, desarrollo, conclusión)")
        suggestions.append("Haz pausas de 15 minutos entre pomodoros")
    elif any(keyword in title_lower or keyword in notes_lower for keyword in medium_keywords):
        estimated = 3
        confidence = 0.8
        reasoning = f"📝 Tarea de complejidad media. Requiere concentración sostenida."
        suggestions.append("Elimina distracciones durante las sesiones")
        suggestions.append("Prepara materiales antes de comenzar")
    elif any(keyword in title_lower or keyword in notes_lower for keyword in short_keywords):
        estimated = 1
        confidence = 0.9
        reasoning = f"⚡ Tarea corta y específica. Una sesión debería ser suficiente."
        suggestions.append("Complétala de una sola vez para mejor retención")
    
    # Ajuste por longitud de notas
    if notes_length > 500:
        estimated += 2
        reasoning += f" Las notas extensas ({notes_length} caracteres) indican complejidad adicional."
    elif notes_length > 250:
        estimated += 1
        reasoning += f" Las notas detalladas sugieren una tarea más compleja."
    
    # Ajuste por curso (algunos cursos son más demandantes)
    difficult_courses = ['matemáticas', 'física', 'química', 'cálculo', 'estadística', 
                        'programación', 'algoritmos', 'estructuras de datos']
    if any(course in course_lower for course in difficult_courses):
        estimated += 1
        confidence *= 0.9
        reasoning += " La asignatura requiere alto nivel de concentración."
        suggestions.append(f"📊 {request.course} requiere sesiones sin interrupciones")
    
    # Ajuste por urgencia
    if days_until_due is not None:
        if days_until_due <= 1:
            suggestions.append("⚠️ ¡Fecha límite muy próxima! Prioriza esta tarea")
        elif days_until_due <= 3:
            suggestions.append("📅 Fecha límite cercana. Planifica tus pomodoros pronto")
    
    # Límites razonables
    estimated = max(1, min(10, estimated))
    
    # Mensajes motivacionales
    if estimated <= 2:
        suggestions.append("💪 ¡Puedes terminarlo rápido! Mantén el enfoque")
    elif estimated >= 6:
        suggestions.append("🎯 Tarea grande. Celebra cada pomodoro completado")
    
    return PomodoroEstimateResponse(
        estimated_pomodoros=estimated,
        confidence=confidence,
        reasoning=reasoning,
        suggestions=suggestions
    )


@app.put("/tasks/{task_id}", response_model=Task)
def update_task(task_id: int, task: Task, authorization: Optional[str] = Header(None)) -> Task:
    """Reemplaza el contenido de una tarea existente."""
    session = _require_session(authorization)
    tasks = load_tasks(session.email)
    for index, existing in enumerate(tasks):
        if existing.id == task_id:
            task.id = task_id
            task.user_email = session.email  # Asignar email del usuario
            tasks[index] = task
            save_tasks(tasks, user_email=session.email)
            return task
    raise HTTPException(status_code=404, detail="Tarea no encontrada")


@app.patch("/tasks/{task_id}/status", response_model=Task)
def update_task_status(task_id: int, status: TaskStatus, authorization: Optional[str] = Header(None)) -> Task:
    """Actualiza solo el estado de una tarea sin modificar el resto de campos.
    
    Acepta status tanto como query parameter como en el body.
    Ejemplo: PATCH /tasks/1/status?status=completed
    """
    session = _require_session(authorization)
    tasks = load_tasks(session.email)
    for index, existing in enumerate(tasks):
        if existing.id == task_id:
            # Asegurarse de que status es un string (puede venir como QueryParam)
            status_value = status.value if isinstance(status, TaskStatus) else str(status)
            # Validar que sea un status válido
            try:
                status_enum = TaskStatus(status_value)
            except ValueError:
                raise HTTPException(status_code=400, detail=f"Estado inválido: {status_value}")
            
            updated = _model_copy(existing, update={"status": status_enum})
            updated.user_email = session.email  # Asegurarse de que tenga el email
            tasks[index] = updated
            save_tasks(tasks, user_email=session.email)
            return updated
    raise HTTPException(status_code=404, detail="Tarea no encontrada")


@app.delete("/tasks/{task_id}", status_code=204)
def delete_task(task_id: int, authorization: Optional[str] = Header(None)) -> None:
    """Elimina una tarea por identificador."""
    session = _require_session(authorization)
    existing_tasks = load_tasks(session.email)
    tasks = [task for task in existing_tasks if task.id != task_id]
    if len(tasks) == len(existing_tasks):
        raise HTTPException(status_code=404, detail="Tarea no encontrada")
    save_tasks(tasks, user_email=session.email)

    # Cancelar notificaciones pendientes de esta tarea
    try:
        from .notification_scheduler import cancel_entity_notifications
        cancel_entity_notifications("task", task_id, session.email)
    except Exception as e:
        logger.debug("No se pudieron cancelar notificaciones de tarea: %s", e)


@app.get("/tasks/{task_id}/estimate-time", response_model=TimeEstimateResponse)
def estimate_task_time(task_id: int, authorization: Optional[str] = Header(None)) -> TimeEstimateResponse:
    """Estima cuánto tiempo (en minutos) tomará completar una tarea específica."""
    session = _require_session(authorization)
    tasks = load_tasks(session.email)
    
    # Encontrar la tarea
    task = None
    for t in tasks:
        if t.id == task_id:
            task = t
            break
    
    if not task:
        raise HTTPException(status_code=404, detail="Tarea no encontrada")
    
    # Analizar el contenido de la tarea
    title_lower = task.title.lower()
    notes_lower = (task.notes or "").lower()
    course_lower = (task.course or "").lower()
    
    # Palabras clave que indican tareas largas
    long_keywords = ['ensayo', 'proyecto', 'investigación', 'tesis', 'monografía', 'informe extenso', 
                     'presentación completa', 'análisis profundo', 'desarrollo completo', 'examen final']
    medium_keywords = ['tarea', 'ejercicios', 'práctica', 'leer capítulo', 'estudio', 'repaso',
                       'laboratorio', 'informe', 'resumen', 'presentación', 'quiz']
    short_keywords = ['test corto', 'lectura breve', 'repasar', 'revisar', 'corregir', 'completar']
    
    # Calcular longitud del contenido
    notes_length = len(task.notes or "")
    
    # Análisis de urgencia
    days_until_due = None
    if task.due_date:
        try:
            days_until_due = (task.due_date - datetime.now(timezone.utc)).days
        except:
            days_until_due = None
    
    # Calcular estimación base (en minutos, convertida de pomodoros)
    estimated_pomodoros = 2  # Por defecto
    confidence = 0.7
    reasoning = ""
    suggestions = []
    
    # Ajuste por palabras clave
    if any(keyword in title_lower or keyword in notes_lower for keyword in long_keywords):
        estimated_pomodoros = 5
        confidence = 0.85
        reasoning = "📚 Tarea de tipo académico extenso detectado. Requiere múltiples sesiones de concentración profunda."
        suggestions.append("Divide la tarea en secciones (introducción, desarrollo, conclusión)")
        suggestions.append("Haz pausas de 15 minutos entre sesiones")
    elif any(keyword in title_lower or keyword in notes_lower for keyword in medium_keywords):
        estimated_pomodoros = 3
        confidence = 0.8
        reasoning = "📝 Tarea de complejidad media. Requiere concentración sostenida."
        suggestions.append("Elimina distracciones durante las sesiones")
        suggestions.append("Prepara materiales antes de comenzar")
    elif any(keyword in title_lower or keyword in notes_lower for keyword in short_keywords):
        estimated_pomodoros = 1
        confidence = 0.9
        reasoning = "⚡ Tarea corta y específica. Una sesión debería ser suficiente."
        suggestions.append("Complétala de una sola vez para mejor retención")
    
    # Ajuste por longitud de notas
    if notes_length > 500:
        estimated_pomodoros += 2
        reasoning += f" Las notas extensas ({notes_length} caracteres) indican complejidad adicional."
    elif notes_length > 250:
        estimated_pomodoros += 1
        reasoning += f" Las notas detalladas sugieren una tarea más compleja."
    
    # Ajuste por curso (algunos cursos son más demandantes)
    difficult_courses = ['matemáticas', 'física', 'química', 'cálculo', 'estadística', 
                        'programación', 'algoritmos', 'estructuras de datos', 'ingeniería']
    if any(course in course_lower for course in difficult_courses):
        estimated_pomodoros += 1
        confidence *= 0.9
        reasoning += " La asignatura requiere alto nivel de concentración."
        suggestions.append(f"📊 {task.course} requiere sesiones sin interrupciones")
    
    # Ajuste por urgencia
    if days_until_due is not None:
        if days_until_due <= 1:
            suggestions.append("⚠️ ¡Fecha límite muy próxima! Prioriza esta tarea")
        elif days_until_due <= 3:
            suggestions.append("📅 Fecha límite cercana. Planifica tus sesiones pronto")
    
    # Límites razonables
    estimated_pomodoros = max(1, min(10, estimated_pomodoros))
    
    # Mensajes motivacionales
    if estimated_pomodoros <= 2:
        suggestions.append("💪 ¡Puedes terminarlo rápido! Mantén el enfoque")
    elif estimated_pomodoros >= 6:
        suggestions.append("🎯 Tarea grande. Celebra cada etapa completada")
    
    # Convertir pomodoros a minutos (25 minutos por pomodoro)
    estimated_minutes = estimated_pomodoros * 25
    
    return TimeEstimateResponse(
        estimated_minutes=estimated_minutes,
        confidence=confidence,
        reasoning=reasoning,
        suggestions=suggestions
    )


# ════════════════════════════════════════════════════════════════════
# RECORDATORIOS (CRUD completo):
#   GET    /reminders        → lista recordatorios del usuario (caché 30 s)
#   POST   /reminders        → crea recordatorio, crea evento en Google/Microsoft
#                              Calendar si hay tokens, y programa push
#   DELETE /reminders/{id}  → elimina recordatorio, evento del calendario
#                              y cancela notificaciones push
#   PATCH  /reminders/{id}  → actualiza campos y reprograma notificaciones
#
# Cada recordatorio tiene: título, descripción, fecha/hora (UTC),
# tipo (TASK/FOCUS/PERSONAL) y opcionalmente un calendar_event_id
# si se sincronizó con Google o Microsoft Calendar.
# ════════════════════════════════════════════════════════════════════
@app.get("/reminders", response_model=list[Reminder])
def list_reminders(authorization: Optional[str] = Header(None)) -> list[Reminder]:
    """Devuelve los recordatorios almacenados."""
    session = _require_session(authorization)
    return _get_cached_reminders(session.email)  # Use cache


@app.post("/reminders", response_model=Reminder, status_code=201)
async def create_reminder(
    reminder: ReminderCreate,
    background_tasks: BackgroundTasks,
    authorization: Optional[str] = Header(None),
    x_timezone: Optional[str] = Header(None, alias="X-Timezone"),
) -> Reminder:
    """Crea un recordatorio y lo programa en el calendario del usuario."""
    session = _require_session(authorization)
    reminders = _get_cached_reminders(session.email)  # Use cache
    reminder_id = next_id(reminders)
    new_reminder = Reminder(
        id=reminder_id,
        title=reminder.title,
        description=reminder.description,
        remind_at=_ensure_utc(reminder.remind_at),
        type=reminder.type,
        delivery_provider=session.provider,
        user_email=session.email
    )

    # Crear evento en Calendar (Google o Microsoft)
    if session.provider in (AuthProvider.GOOGLE, AuthProvider.MICROSOFT):
        try:
            tokens_map = _get_cached_tokens()  # Use cache
            tokens = tokens_map.get(session.email)
            if tokens and tokens.get("access_token"):
                event_id = await calendar_integration.create_calendar_event(
                    new_reminder, session, tokens, user_email=session.email
                )
                if event_id:
                    new_reminder.calendar_event_id = event_id
            else:
                logger.info("No se encontraron tokens válidos para %s; omitiendo creación en Calendar.", session.email)
        except Exception as e:
            logger.warning("No se pudo crear el evento en Calendar: %s", e)

    reminders.append(new_reminder)
    from .storage import save_reminders as _save_reminders
    _save_reminders(reminders, user_email=session.email)
    _invalidate_reminders_cache(session.email)  # Invalidate cache

    # Programar notificaciones push en background (no bloquea la respuesta)
    _new_reminder_ref = new_reminder  # capturar referencia para el closure
    _user_email = session.email
    def _schedule_reminder_notifs():
        try:
            from .notification_scheduler import schedule_reminder_notifications
            schedule_reminder_notifications(
                reminder_id=_new_reminder_ref.id,
                title=_new_reminder_ref.title,
                description=_new_reminder_ref.description,
                remind_at=_new_reminder_ref.remind_at,
                user_email=_user_email,
                user_timezone=x_timezone,
            )
        except Exception as e:
            logger.warning("⚠️ No se pudieron programar notificaciones de recordatorio: %s", e)

    background_tasks.add_task(_schedule_reminder_notifs)

    # Enviar email del recordatorio (sin push, el scheduler se encarga)
    _send_reminder_notification(new_reminder)

    # Ensure we're returning the reminder with the correct ID
    logger.info(f"✅ Created reminder with ID {new_reminder.id} for {session.email}")
    return new_reminder


@app.delete("/reminders/{reminder_id}", status_code=204)
async def delete_reminder(reminder_id: str, authorization: Optional[str] = Header(None)) -> None:
    """Elimina un recordatorio existente y su evento de calendario asociado."""
    session = _require_session(authorization)
    reminders = _get_cached_reminders(session.email)  # Use cache
    
    # Convertir a int si es posible
    try:
        reminder_id_value = int(reminder_id)
    except ValueError:
        # Si no se puede convertir a int, es un ID temporal - no existe en BD
        logger.warning(f"Intento de eliminar recordatorio con ID temporal: {reminder_id}")
        raise HTTPException(status_code=404, detail="Recordatorio no encontrado (ID temporal)")
    
    # Encontrar el recordatorio a eliminar
    reminder_to_delete = None
    for reminder in reminders:
        if reminder.id == reminder_id_value:
            reminder_to_delete = reminder
            break
    
    if not reminder_to_delete:
        logger.warning(f"Recordatorio {reminder_id} no encontrado para usuario {session.email}")
        raise HTTPException(status_code=404, detail="Recordatorio no encontrado")
    
    logger.info(f"🗑️ Eliminando recordatorio {reminder_id_value} para {session.email}")
    
    # Eliminar evento del calendario si existe
    if reminder_to_delete.calendar_event_id:
        try:
            tokens_map = _get_cached_tokens()  # Use cache
            tokens = tokens_map.get(session.email)
            if tokens:
                await calendar_integration.delete_calendar_event(reminder_to_delete, tokens, user_email=session.email)
                logger.info(f"✅ Evento de calendario eliminado para recordatorio {reminder_id}")
            else:
                logger.warning(f"No se encontraron tokens para eliminar evento del recordatorio {reminder_id}")
        except Exception as e:
            logger.warning(f"No se pudo eliminar el evento del calendario: {e}")
            # Continuamos eliminando el recordatorio aunque falle eliminar del calendario
    
    # Eliminar el recordatorio de la base de datos
    updated = [reminder for reminder in reminders if reminder.id != reminder_id_value]
    from .storage import save_reminders as _save_reminders
    _save_reminders(updated, user_email=session.email)
    _invalidate_reminders_cache(session.email)  # Invalidate cache

    # Cancelar notificaciones pendientes de este recordatorio
    try:
        from .notification_scheduler import cancel_entity_notifications
        cancel_entity_notifications("reminder", reminder_id_value, session.email)
    except Exception as e:
        logger.debug("No se pudieron cancelar notificaciones de recordatorio: %s", e)


@app.patch("/reminders/{reminder_id}", response_model=Reminder)
def update_reminder(
    reminder_id: int,
    update: ReminderUpdate,
    background_tasks: BackgroundTasks,
    authorization: Optional[str] = Header(None),
    x_timezone: Optional[str] = Header(None, alias="X-Timezone"),
) -> Reminder:
    """Permite ajustar título, notas, hora o tipo de un recordatorio."""

    session = _require_session(authorization)
    reminders = _get_cached_reminders(session.email)  # Use cache
    for index, existing in enumerate(reminders):
        if existing.id != reminder_id:
            continue

        changes: dict[str, object] = {}
        if update.title is not None:
            changes["title"] = update.title
        if update.description is not None:
            changes["description"] = update.description
        if update.remind_at is not None:
            changes["remind_at"] = _ensure_utc(update.remind_at)
        if update.type is not None:
            changes["type"] = update.type

        if not changes:
            return existing

        updated_reminder = _model_copy(
            existing,
            update={
                **changes,
                "delivery_provider": session.provider,
            },
        )
        reminders[index] = updated_reminder
        from .storage import save_reminders as _save_reminders
        _save_reminders(reminders, user_email=session.email)
        _invalidate_reminders_cache(session.email)  # Invalidate cache

        # Reprogramar notificaciones push en background (no bloquea la respuesta)
        _updated = updated_reminder
        _rid = reminder_id
        _email = session.email
        _tz = x_timezone
        def _reschedule_reminder_notifs():
            try:
                from .notification_scheduler import cancel_entity_notifications, schedule_reminder_notifications
                cancel_entity_notifications("reminder", _rid, _email)
                schedule_reminder_notifications(
                    reminder_id=_updated.id,
                    title=_updated.title,
                    description=_updated.description,
                    remind_at=_updated.remind_at,
                    user_email=_email,
                    user_timezone=_tz,
                )
            except Exception as e:
                logger.debug("No se pudieron reprogramar notificaciones: %s", e)

        background_tasks.add_task(_reschedule_reminder_notifs)

        _send_reminder_notification(updated_reminder)
        return updated_reminder

    raise HTTPException(status_code=404, detail="Recordatorio no encontrado")


# ════════════════════════════════════════════════════════════════════
# HORARIO SEMANAL (schedule):
#   GET    /schedule        → devuelve todos los bloques del horario del usuario
#   POST   /schedule        → añade un bloque de clase/actividad;
#                             valida que end_time > start_time
#   DELETE /schedule/{id}  → elimina un bloque directamente de Supabase
#
# Cada entrada tiene: día de la semana, hora inicio/fin, materia,
# aula y color. Persiste por usuario en la tabla schedule_entries.
# ════════════════════════════════════════════════════════════════════
@app.get("/schedule", response_model=list[ScheduleEntry])
def list_schedule(authorization: Optional[str] = Header(None)) -> list[ScheduleEntry]:
    """Entrega todos los bloques del horario semanal del usuario autenticado."""
    session = _require_session(authorization)
    return load_schedule(user_email=session.email)


@app.post("/schedule", response_model=ScheduleEntry, status_code=201)
def create_schedule_entry(entry: ScheduleEntry, authorization: Optional[str] = Header(None)) -> ScheduleEntry:
    """Registra un bloque de horario validando que la hora final sea posterior."""
    session = _require_session(authorization)
    if entry.end_time <= entry.start_time:
        raise HTTPException(status_code=400, detail="La hora de término debe ser posterior al inicio.")

    entries = load_schedule(user_email=session.email)
    entry.id = next_id(entries)
    entries.append(entry)
    
    # Guardar con el email del usuario autenticado
    from .storage import save_schedule as _save_schedule
    _save_schedule(entries, user_email=session.email)
    return entry


@app.delete("/schedule/{entry_id}", status_code=204)
def delete_schedule_entry(entry_id: int, authorization: Optional[str] = Header(None)) -> None:
    """Elimina un bloque del horario semanal del usuario autenticado."""
    session = _require_session(authorization)
    
    try:
        supabase = get_supabase()
        # Delete directly from the database
        supabase.table("schedule_entries").delete().eq("id", entry_id).eq("user_email", session.email).execute()
    except Exception as e:
        logger.warning(f"Error eliminando horario: {e}")
        raise HTTPException(status_code=500, detail="No se pudo eliminar el bloque del horario")


# ════════════════════════════════════════════════════════════════════
# SESIONES DE ENFOQUE (Pomodoro/Focus):
#   GET  /focus-sessions → historial de todas las sesiones de estudio
#   POST /focus-sessions → registra una nueva sesión con duración y tipo
#
# El frontend usa estos datos para calcular estadísticas de
# productividad (tiempo total estudiado, rachas, logros).
# NOTA: estos endpoints no validan sesión; se puede restringir
# en el futuro agregando _require_session().
# ════════════════════════════════════════════════════════════════════
@app.get("/focus-sessions", response_model=list[FocusSession])
def list_focus_sessions(
    authorization: Optional[str] = Header(None)
) -> list[FocusSession]:
    """Recupera el historial de sesiones de enfoque (pomodoro) del usuario autenticado."""
    # Validar autenticación
    if not authorization or not authorization.startswith("Bearer "):
        raise HTTPException(status_code=401, detail="Missing or invalid authorization")
    
    session_token = authorization[7:]
    session_data = validate_session_token(session_token)
    if not session_data:
        raise HTTPException(status_code=401, detail="Invalid session token")
    
    user_email = session_data.get("email")
    return load_focus_sessions_for_user(user_email)


@app.post("/focus-sessions", response_model=FocusSession, status_code=201)
def create_focus_session(
    session: FocusSession,
    authorization: Optional[str] = Header(None)
) -> FocusSession:
    """Guarda una nueva sesión de enfoque con duración y fecha."""
    # Validar autenticación
    if not authorization or not authorization.startswith("Bearer "):
        raise HTTPException(status_code=401, detail="Missing or invalid authorization")
    
    session_token = authorization[7:]
    session_data = validate_session_token(session_token)
    if not session_data:
        raise HTTPException(status_code=401, detail="Invalid session token")
    
    user_email = session_data.get("email")
    
    # Cargar SOLO las sesiones del usuario actual para calcular el ID
    user_sessions = load_focus_sessions_for_user(user_email)
    session.id = next_id(user_sessions)
    session.user_email = user_email  # Asociar al usuario autenticado
    
    # Guardar la nueva sesión en Supabase
    try:
        supabase = get_supabase()
        data = {
            "id": session.id,
            "user_email": user_email,
            "topic": session.topic,
            "duration_minutes": session.duration_minutes,
            "completed_at": session.completed_at.isoformat() if session.completed_at else None,
        }
        supabase.table("focus_sessions").insert(data).execute()
        logger.info(f"✅ Focus session guardada para {user_email}: {session.duration_minutes}min ({session.topic})")
    except Exception as e:
        logger.error(f"❌ Error guardando focus session: {e}")
        raise HTTPException(status_code=500, detail=f"Error guardando sesión: {str(e)}")
    
    return session


# ════════════════════════════════════════════════════════════════════
# GOOGLE CALENDAR INTEGRATION:
#   GET    /calendar/google/events         → obtiene eventos del calendario
#                                             del usuario (con refresh automático
#                                             de access_token si expiró)
#   POST   /calendar/google/events         → crea un evento en el calendario
#   DELETE /calendar/google/events/{id}   → elimina un evento
#
# Todos los endpoints requieren que el usuario haya conectado su
# Google Calendar previamente via /auth/google/calendar/start.
# Los tokens se guardan en Supabase con save_token_for_email().
# ════════════════════════════════════════════════════════════════════
# ============================================================================
# Endpoints para Google Calendar Integration
# ============================================================================

@app.get("/calendar/google/events")
async def get_google_calendar_events(
    authorization: Optional[str] = Header(None),
    time_min: Optional[str] = None,
    time_max: Optional[str] = None,
    max_results: int = 50
):
    """Obtiene eventos del Google Calendar del usuario autenticado.
    
    Args:
        time_min: Fecha mínima en formato ISO 8601 (default: hoy)
        time_max: Fecha máxima en formato ISO 8601 (default: 30 días desde hoy)
        max_results: Número máximo de eventos a retornar (default: 50)
    
    Returns:
        Lista de eventos del calendario en formato simplificado
    """
    session = _require_session(authorization)
    logger.info(f"[GOOGLE_CALENDAR] Fetching events for user: {session.email}")
    
    # Obtener tokens de OAuth del usuario
    tokens = load_tokens()
    user_token = tokens.get(session.email)
    
    if not user_token:
        logger.warning(f"[GOOGLE_CALENDAR] No tokens found for {session.email}. Available emails: {list(tokens.keys())}")
        raise HTTPException(
            status_code=401,
            detail="No se encontraron credenciales de Google Calendar. Por favor, vuelve a autenticarte con Google."
        )
    
    access_token = user_token.get("access_token")
    if not access_token:
        logger.warning(f"[GOOGLE_CALENDAR] Token exists but no access_token for {session.email}")
        raise HTTPException(
            status_code=401,
            detail="Token de acceso inválido. Por favor, vuelve a autenticarte con Google."
        )
    
    # Crear cliente de Google Calendar
    calendar_client = await google_calendar.get_calendar_client(access_token)
    
    # Parsear fechas si se proporcionaron
    time_min_dt = None
    time_max_dt = None
    
    if time_min:
        try:
            time_min_dt = datetime.fromisoformat(time_min.replace('Z', '+00:00'))
        except ValueError:
            raise HTTPException(status_code=400, detail="Formato de time_min inválido")
    
    if time_max:
        try:
            time_max_dt = datetime.fromisoformat(time_max.replace('Z', '+00:00'))
        except ValueError:
            raise HTTPException(status_code=400, detail="Formato de time_max inválido")
    
    # Obtener eventos
    try:
        events = await calendar_client.get_events(
            time_min=time_min_dt,
            time_max=time_max_dt,
            max_results=max_results
        )
        
        logger.info(f"📅 Obtenidos {len(events)} eventos de Google Calendar para {session.email}")
        return {"events": events, "count": len(events)}
        
    except HTTPException as he:
        # Si es 401, intentar refrescar el token
        if he.status_code == 401 and user_token.get("refresh_token"):
            logger.info(f"[GOOGLE_CALENDAR] Access token expirado, refrescando para {session.email}...")
            try:
                # Refrescar token usando el cliente OAuth
                google_client = get_google_client()
                async with httpx.AsyncClient() as http_client:
                    response = await http_client.post(
                        "https://oauth2.googleapis.com/token",
                        data={
                            "client_id": google_client.config.client_id,
                            "client_secret": google_client.config.client_secret,
                            "refresh_token": user_token["refresh_token"],
                            "grant_type": "refresh_token"
                        }
                    )
                    
                    if response.status_code == 200:
                        new_tokens = response.json()
                        updated_tokens = {
                            **user_token,
                            "access_token": new_tokens['access_token'],
                            "expires_in": new_tokens.get('expires_in', 3599),
                        }
                        save_token_for_email(session.email, updated_tokens)
                        logger.info(f"✅ Token refrescado exitosamente para {session.email}")
                        
                        # Reintentar con nuevo token
                        calendar_client = await google_calendar.get_calendar_client(new_tokens['access_token'])
                        events = await calendar_client.get_events(
                            time_min=time_min_dt,
                            time_max=time_max_dt,
                            max_results=max_results
                        )
                        return {"events": events, "count": len(events)}
                    else:
                        logger.error(f"Error al refrescar token: {response.status_code} - {response.text}")
                        raise HTTPException(
                            status_code=401,
                            detail="Tu sesión de Google Calendar ha expirado. Por favor, vuelve a autenticarte."
                        )
            except Exception as refresh_error:
                logger.error(f"Error al refrescar token: {refresh_error}")
                raise HTTPException(
                    status_code=401,
                    detail="Tu sesión de Google Calendar ha expirado. Por favor, vuelve a autenticarte."
                )
        raise
    except Exception as e:
        logger.error(f"Error inesperado al obtener eventos de Google Calendar: {e}")
        raise HTTPException(
            status_code=500,
            detail="Error al obtener eventos del calendario"
        )


@app.post("/calendar/google/events")
async def create_google_calendar_event(
    authorization: Optional[str] = Header(None),
    summary: str = Form(...),
    start: str = Form(...),
    end: str = Form(...),
    description: Optional[str] = Form(None),
    location: Optional[str] = Form(None)
):
    """Crea un nuevo evento en Google Calendar del usuario.
    
    Args:
        summary: Título del evento
        start: Fecha/hora de inicio en formato ISO 8601
        end: Fecha/hora de fin en formato ISO 8601
        description: Descripción opcional del evento
        location: Ubicación opcional
    
    Returns:
        Evento creado en formato simplificado
    """
    session = _require_session(authorization)
    
    # Obtener tokens de OAuth del usuario
    tokens = load_tokens()
    user_token = tokens.get(session.email)
    
    if not user_token:
        raise HTTPException(
            status_code=401,
            detail="No se encontraron credenciales de Google Calendar"
        )
    
    access_token = user_token.get("access_token")
    if not access_token:
        raise HTTPException(
            status_code=401,
            detail="Token de acceso inválido"
        )
    
    # Parsear fechas
    try:
        start_dt = datetime.fromisoformat(start.replace('Z', '+00:00'))
        end_dt = datetime.fromisoformat(end.replace('Z', '+00:00'))
    except ValueError as e:
        raise HTTPException(
            status_code=400,
            detail=f"Formato de fecha inválido: {e}"
        )
    
    # Crear cliente de Google Calendar
    calendar_client = await google_calendar.get_calendar_client(access_token)
    
    # Crear evento
    try:
        event = await calendar_client.create_event(
            summary=summary,
            start=start_dt,
            end=end_dt,
            description=description,
            location=location
        )
        
        logger.info(f"✅ Evento creado en Google Calendar para {session.email}: {summary}")
        return {"event": event, "success": True}
        
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error inesperado al crear evento: {e}")
        raise HTTPException(
            status_code=500,
            detail="Error al crear evento en el calendario"
        )


@app.delete("/calendar/google/events/{event_id}")
async def delete_google_calendar_event(
    event_id: str,
    authorization: Optional[str] = Header(None)
):
    """Elimina un evento del Google Calendar del usuario.
    
    Args:
        event_id: ID del evento en Google Calendar (similar a "7jbmf8h6g7h8j9k0l")
    
    Returns:
        Confirmación de eliminación exitosa
    """
    session = _require_session(authorization)
    logger.info(f"[GOOGLE_CALENDAR] Deleting event {event_id} for user: {session.email}")
    
    # Obtener tokens de OAuth del usuario
    tokens = load_tokens()
    user_token = tokens.get(session.email)
    
    if not user_token:
        logger.warning(f"[GOOGLE_CALENDAR] No tokens found for {session.email}")
        raise HTTPException(
            status_code=401,
            detail="No se encontraron credenciales de Google Calendar. Por favor, vuelve a autenticarte con Google."
        )
    
    access_token = user_token.get("access_token")
    if not access_token:
        logger.warning(f"[GOOGLE_CALENDAR] Token exists but no access_token for {session.email}")
        raise HTTPException(
            status_code=401,
            detail="Token de acceso inválido. Por favor, vuelve a autenticarte con Google."
        )
    
    # Crear cliente de Google Calendar
    calendar_client = await google_calendar.get_calendar_client(access_token)
    
    # Eliminar evento
    try:
        await calendar_client.delete_event(event_id)
        logger.info(f"✅ Evento {event_id} eliminado de Google Calendar para {session.email}")
        return {"success": True, "message": f"Evento {event_id} eliminado correctamente"}
        
    except HTTPException as he:
        # Si es 401, intentar refrescar el token
        if he.status_code == 401 and user_token.get("refresh_token"):
            logger.info(f"[GOOGLE_CALENDAR] Access token expirado, refrescando para {session.email}...")
            try:
                # Refrescar token usando el cliente OAuth
                google_client = get_google_client()
                async with httpx.AsyncClient() as http_client:
                    response = await http_client.post(
                        "https://oauth2.googleapis.com/token",
                        data={
                            "client_id": google_client.config.client_id,
                            "client_secret": google_client.config.client_secret,
                            "refresh_token": user_token["refresh_token"],
                            "grant_type": "refresh_token"
                        }
                    )
                    
                    if response.status_code == 200:
                        new_tokens = response.json()
                        updated_tokens = {
                            **user_token,
                            "access_token": new_tokens['access_token'],
                            "expires_in": new_tokens.get('expires_in', 3599),
                        }
                        save_token_for_email(session.email, updated_tokens)
                        logger.info(f"✅ Token refrescado exitosamente para {session.email}")
                        
                        # Reintentar con nuevo token
                        calendar_client = await google_calendar.get_calendar_client(new_tokens['access_token'])
                        await calendar_client.delete_event(event_id)
                        logger.info(f"✅ Evento {event_id} eliminado después de refrescar token")
                        return {"success": True, "message": f"Evento {event_id} eliminado correctamente"}
                    else:
                        logger.error(f"Error al refrescar token: {response.status_code} - {response.text}")
                        raise HTTPException(
                            status_code=401,
                            detail="Tu sesión de Google Calendar ha expirado. Por favor, vuelve a autenticarte."
                        )
            except Exception as refresh_error:
                logger.error(f"Error al refrescar token: {refresh_error}")
                raise HTTPException(
                    status_code=401,
                    detail="Tu sesión de Google Calendar ha expirado. Por favor, vuelve a autenticarte."
                )
        raise
    except Exception as e:
        logger.error(f"Error inesperado al eliminar evento: {e}")
        raise HTTPException(
            status_code=500,
            detail="Error al eliminar evento del calendario"
        )


# ============================================================================
# Endpoints para sincronizar estadísticas de gamificación con la base de datos
# ============================================================================

@app.post("/user-stats")
def save_user_stats_endpoint(
    payload: UserStatsPayload,
    authorization: Optional[str] = Header(None)
) -> dict:
    """Guarda las estadísticas de gamificación del usuario en la base de datos."""
    session = _require_session(authorization)
    
    from .supabase_storage import save_user_stats
    save_user_stats(
        user_email=session.email,
        xp=payload.xp,
        streak_days=payload.streak_days,
        last_activity_date=payload.last_activity_date,
        total_tasks_ever_completed=payload.total_tasks_ever_completed,
        unlocked_achievements=payload.unlocked_achievements
    )
    
    logger.info(f"✅ Saved user stats for {session.email}: XP={payload.xp}, Streak={payload.streak_days}, Achievements={len(payload.unlocked_achievements)}")
    return {
        "success": True,
        "email": session.email,
        "message": "User stats saved successfully"
    }


@app.get("/user-stats")
def load_user_stats_endpoint(authorization: Optional[str] = Header(None)) -> dict:
    """Carga las estadísticas de gamificación del usuario desde la base de datos."""
    session = _require_session(authorization)
    
    from .supabase_storage import load_user_stats
    stats = load_user_stats(session.email)
    
    logger.info(f"✅ Loaded user stats for {session.email}")
    return stats


# ════════════════════════════════════════════════════════════════════
# DASHBOARD Y ESTADÍSTICAS DE GAMIFICACIÓN:
#   POST /user-stats          → guarda o actualiza XP, racha y nivel del usuario
#   GET  /user-stats          → devuelve las estadísticas actuales del usuario
#   GET  /dashboard           → computa estadísticas rápidas (tareas, rachas,
#                                sesiones) para el tablero principal
#   POST /stats               → guarda explícitamente un objeto DashboardStats
#   PUT  /api/update-streak   → actualiza la racha de días del usuario
#
# La gamificación persiste en la tabla user_stats de Supabase:
# streak_days, total_xp, level, achievements (JSON).
# ════════════════════════════════════════════════════════════════════
@app.get("/dashboard", response_model=DashboardStats)
def dashboard(authorization: Optional[str] = Header(None)) -> DashboardStats:
    """Calcula estadísticas rápidas para el tablero principal."""
    session = _require_session(authorization)
    return compute_dashboard_stats(load_tasks(session.email), load_reminders(session.email), load_focus_sessions())


# ════════════════════════════════════════════════════════════════════
# RESUMEN E IA (NotebookLM / Google Gemini):
#   POST /summary          → resumen NLP local desde texto o archivo subido
#   POST /summary/text     → resumen NLP solo desde texto (sin archivo)
#   POST /summary/ai       → resumen mejorado con Google Gemini; si la IA no
#                             está disponible hace fallback al algoritmo local.
#                             Soporta extensiones: short / medium / long
#   GET  /summary/ai/status→ indica si la IA de resumen está activa
#   GET  /ai/stats         → estadísticas de uso de la IA (total llamadas, tokens)
#   POST /ai/sentiment     → análisis de sentimiento de un texto
#   POST /ai/categorize    → clasifica texto en categorías predefinidas
#   POST /ai/entities      → extrae entidades (personas, lugares, fechas)
#   POST /ai/generate      → generación libre de contenido a partir de un prompt
#   POST /ai/qa            → responde preguntas sobre un documento
#   POST /ai/analyze-image → análisis visual de una imagen subida
#   POST /ai/generate-image→ genera una imagen a partir de un prompt
#   POST /ai/extract-text  → extrae texto de TXT, PDF, DOCX, PPTX
# ════════════════════════════════════════════════════════════════════
@app.post("/summary", response_model=SummaryResponse)
async def create_summary(
    text: Optional[str] = Form(None),
    sentences: int = Form(5),
    file: Optional[UploadFile] = File(None),
) -> SummaryResponse:
    """Genera un resumen desde un archivo o texto crudo y extrae palabras clave."""
    try:
        logger.info(f"📄 Summary request - file: {file.filename if file else 'None'}, text length: {len(text) if text else 0}")
        summary_text, original_text = await summarizer.generate_summary(file, text, sentences)
        source_for_keywords = original_text or text or summary_text
        keywords = summarizer.top_keywords(source_for_keywords)
        global _summaries_generated
        _summaries_generated += 1
        return SummaryResponse(
            summary=summary_text,
            highlighted_keywords=keywords,
            original_text=original_text,
        )
    except Exception as e:
        logger.error(f"❌ Error generating summary: {e}")
        import traceback
        logger.error(traceback.format_exc())
        raise HTTPException(status_code=500, detail=f"Error generating summary: {str(e)}")


@app.post("/summary/text", response_model=SummaryResponse)
async def create_summary_from_text(request: SummaryRequest) -> SummaryResponse:
    """Permite generar resúmenes exclusivamente desde texto enviado en JSON."""
    summary_text, original_text = await summarizer.summarize_from_text(
        request.text or "",
        request.sentences,
    )
    keywords = summarizer.top_keywords(original_text or request.text or summary_text)
    global _summaries_generated
    _summaries_generated += 1
    notebooklm_tasks.notebooklm_tasks.summaries_count += 1
    try:
        from .supabase_storage import increment_summaries_count
        increment_summaries_count()
    except Exception:
        pass
    return SummaryResponse(
        summary=summary_text,
        highlighted_keywords=keywords,
        original_text=original_text,
    )


@app.post("/summary/ai", response_model=SummaryResponse)
async def create_summary_with_ai(
    text: Optional[str] = Form(None),
    sentences: int = Form(5),
    file: Optional[UploadFile] = File(None),
    use_ai: bool = Form(True),
    summary_length: str = Form("medium"),  # "short", "medium", "long"
) -> SummaryResponse:
    """
    Genera un resumen usando NotebookLM AI (alternativa más inteligente).
    Si NotebookLM no está disponible, usa el algoritmo local automáticamente.
    
    Parámetros:
        text: Texto a resumir (opcional)
        sentences: Número de frases en el resumen
        file: Archivo a subir (opcional)
        use_ai: Usar IA (default: True, fallback automático)
        summary_length: Extensión del resumen
                       - "short": Resumen muy breve (1-3 frases, rápido)
                       - "medium": Resumen equilibrado (5-10 frases, recomendado)
                       - "long": Resumen detallado (10+ frases, comprensivo)
    
    Nota: Esta es una alternativa mejorada a /summary que usa IA cuando está disponible.
    """
    # Valida summary_length
    if summary_length not in ("short", "medium", "long"):
        summary_length = "medium"
    
    # Extrae el texto del archivo si se proporciona
    if file is not None:
        extracted_text, _ = await summarizer.generate_summary(file, None, 1)
    else:
        extracted_text = text or ""
    
    if not extracted_text.strip():
        raise HTTPException(status_code=400, detail="No text provided")
    
    # Intenta usar IA, fallback a algoritmo local
    try:
        if use_ai:
            summary_text, _ = await notebooklm_integration.get_summary_with_ai(
                extracted_text, 
                sentences, 
                prefer_ai=True,
                summary_length=summary_length  # Pasar extensión a la IA
            )
            keywords = await notebooklm_integration.get_keywords_with_ai(
                extracted_text, k=10, prefer_ai=True
            )
        else:
            summary_text, _ = await summarizer.summarize_from_text(extracted_text, sentences)
            keywords = summarizer.top_keywords(extracted_text)
    except Exception as e:
        logging.error(f"Error en resumen con IA: {e}")
        # Fallback: usa algoritmo local
        summary_text, _ = await summarizer.summarize_from_text(extracted_text, sentences)
        keywords = summarizer.top_keywords(extracted_text)
    
    global _summaries_generated
    _summaries_generated += 1
    notebooklm_tasks.notebooklm_tasks.summaries_count += 1
    logger.warning(f"📝 SUMMARY generado — intentando persistir en Supabase (_mem={_summaries_generated})")
    try:
        from .supabase_storage import increment_summaries_count
        result = increment_summaries_count()
        logger.warning(f"📝 increment_summaries_count retornó: {result}")
    except Exception as e:
        logger.warning(f"📝 EXCEPCIÓN en increment_summaries_count: {e}")
    return SummaryResponse(
        summary=summary_text,
        highlighted_keywords=keywords,
        original_text=extracted_text,
    )


@app.get("/summary/ai/status")
async def get_ai_status() -> dict:
    """
    Verifica el estado de la integración con NotebookLM AI (Google Gemini).
    
    Returns:
        {
            "ai_available": True/False,
            "service": "NotebookLM",
            "models": ["gemini-1.5-flash"],
            "message": "IA disponible / IA no configurada / Fallback a algoritmo local"
        }
    """
    is_available = notebooklm_integration.notebooklm_summarizer.is_available()
    
    if is_available:
        status = "IA disponible - Resúmenes potenciados con NotebookLM (Google Gemini)"
    else:
        status = "IA no disponible - Se usa algoritmo local"
    
    return {
        "ai_available": is_available,
        "service": "NotebookLM" if is_available else "Local Algorithm",
        "models": ["gemini-1.5-flash"],
        "message": status,
    }


# ======================== ENDPOINTS DE TAREAS AVANZADAS CON NOTEBOOKLM ========================

@app.get("/ai/stats")
async def get_ai_stats() -> dict:
    """
    Obtiene estadísticas de uso de la API de NotebookLM/Gemini.
    
    Returns:
        {
            "requests_today": int,      # Requests realizados hoy
            "cache_hits": int,          # Respuestas del caché
            "requests_remaining": int,  # Requests restantes (~20/día gratis)
            "cached_responses": int     # Total de respuestas en caché
        }
    """
    try:
        from .notebooklm_tasks import notebooklm_tasks
        stats = notebooklm_tasks.get_stats()
        return {
            "success": True,
            "stats": stats,
            "message": f"✅ Quedan ~{stats['requests_remaining']} requests. Cache hits: {stats['cache_hits']}"
        }
    except Exception as e:
        logging.error(f"❌ Error obteniendo stats: {str(e)}")
        return {"error": f"Error: {str(e)}"}

@app.post("/ai/sentiment")
async def sentiment_analysis(text: str = Form(...)) -> dict:
    """
    Analiza el sentimiento del texto usando NotebookLM (Google Gemini).
    
    Args:
        text: Texto a analizar
    
    Returns:
        {
            "sentiment": "Positivo|Neutral|Negativo",
            "score": 0.0-1.0,
            "explanation": "Explicación breve"
        }
    """
    try:
        if not text or not text.strip():
            return {"error": "El texto no puede estar vacío"}
        return await notebooklm_tasks.analyze_sentiment(text)
    except Exception as e:
        logging.error(f"❌ Error en /ai/sentiment: {str(e)}")
        return {"error": f"Error en análisis de sentimiento: {str(e)}"}


@app.post("/ai/categorize")
async def categorize_text(
    text: str = Form(...),
    categories: Optional[str] = Form(None)
) -> dict:
    """
    Categoriza el texto automáticamente.
    
    Args:
        text: Texto a categorizar
        categories: Categorías separadas por coma (opcional)
    
    Returns:
        {
            "category": "Categoría detectada",
            "confidence": 0.0-1.0,
            "alternatives": [...]
        }
    """
    try:
        if not text or not text.strip():
            return {"error": "El texto no puede estar vacío"}
        cats = None
        if categories:
            cats = [c.strip() for c in categories.split(",")]
        return await notebooklm_tasks.categorize(text, cats)
    except Exception as e:
        logging.error(f"❌ Error en /ai/categorize: {str(e)}")
        return {"error": f"Error en categorización: {str(e)}"}


@app.post("/ai/entities")
async def extract_named_entities(text: str = Form(...)) -> dict:
    """
    Extrae entidades nombradas (NER) del texto usando NotebookLM (Google Gemini).
    
    Returns:
        {
            "personas": [...],
            "lugares": [...],
            "organizaciones": [...]
        }
    """
    try:
        if not text or not text.strip():
            return {"error": "El texto no puede estar vacío"}
        return await notebooklm_tasks.extract_entities(text)
    except Exception as e:
        logging.error(f"❌ Error en /ai/entities: {str(e)}")
        return {"error": f"Error en extracción de entidades: {str(e)}"}


@app.post("/ai/generate")
async def generate_with_ai(
    prompt: str = Form(...),
    max_length: int = Form(8000)
) -> dict:
    """
    Genera contenido basado en un prompt.
    
    Args:
        prompt: Instrucción para generar contenido
        max_length: Longitud máxima en caracteres
    
    Returns:
        {
            "content": "Contenido generado",
            "tokens_used": estimación
        }
    """
    try:
        if not prompt or not prompt.strip():
            return {"error": "El prompt no puede estar vacío"}
        # Detectar si la solicitud es un resumen para contabilizarlo
        _summary_keywords = ('resumen', 'resume', 'sumario', 'summary', 'resumir', 'genera un resumen')
        if any(kw in prompt.lower() for kw in _summary_keywords):
            notebooklm_tasks.notebooklm_tasks.summaries_count += 1
            try:
                from .supabase_storage import increment_summaries_count
                increment_summaries_count()
            except Exception:
                pass
        return await notebooklm_tasks.generate_content(prompt, max_length)
    except Exception as e:
        logging.error(f"❌ Error en /ai/generate: {str(e)}")
        return {"error": f"Error en generación de contenido: {str(e)}"}


@app.post("/ai/qa")
async def question_answering(
    document: str = Form(...),
    question: str = Form(...)
) -> dict:
    """
    Responde una pregunta basada en un documento usando NotebookLM (Google Gemini).
    
    Args:
        document: Documento de referencia
        question: Pregunta a responder
    
    Returns:
        {
            "answer": "Respuesta",
            "confidence": 0.0-1.0,
            "source_quotes": [...]
        }
    """
    try:
        if not document or not document.strip():
            return {"error": "El documento no puede estar vacío"}
        if not question or not question.strip():
            return {"error": "La pregunta no puede estar vacía"}
        return await notebooklm_tasks.answer_question(document, question)
    except Exception as e:
        logging.error(f"❌ Error en /ai/qa: {str(e)}")
        return {"error": f"Error en Q&A: {str(e)}"}


@app.post("/ai/agent")
async def iris_agent_endpoint(
    message: str = Form(...),
    history: str = Form(default="[]"),
    context: str = Form(default=""),
) -> dict:
    """
    Iris en MODO AGENTE: decide qué herramientas usar (function-calling sobre
    Gemini) y devuelve las acciones a ejecutar más una respuesta natural.

    Returns: { "reply": str, "actions": [{type, args}], "source": str }
    """
    try:
        if not message or not message.strip():
            return {"error": "El mensaje no puede estar vacío", "reply": "", "actions": []}

        import json as _json
        try:
            hist = _json.loads(history) if history else []
            if not isinstance(hist, list):
                hist = []
        except Exception:
            hist = []

        from . import iris_agent

        if not iris_agent.is_available():
            res = await notebooklm_tasks.generate_content(message)
            return {"reply": res.get("content", ""), "actions": [], "source": "fallback"}

        result = iris_agent.run_agent(message.strip(), hist, context)

        # Si el agente falló o no produjo nada, recurre a la generación normal.
        if not result.get("reply") and not result.get("actions"):
            res = await notebooklm_tasks.generate_content(message)
            return {"reply": res.get("content", ""), "actions": [], "source": "fallback"}

        return result
    except Exception as e:
        logging.error(f"❌ Error en /ai/agent: {str(e)}")
        return {"error": f"Error en agente: {str(e)}", "reply": "", "actions": []}


# ─── Análisis multimodal de imágenes ───

@app.post("/ai/analyze-image")
async def analyze_image_endpoint(
    file: UploadFile = File(...),
    prompt: str = Form(default="")
) -> dict:
    """
    Analiza una imagen con Gemini multimodal.
    Extrae texto, describe contenido, identifica conceptos.

    Args:
        file: Imagen (PNG, JPG, WEBP, GIF)
        prompt: Instrucción adicional (opcional)

    Returns:
        { "content": str, "source": "ai" } o { "error": str }
    """
    try:
        image_bytes = await file.read()
        mime_type = file.content_type or "image/png"

        if not image_bytes:
            return {"error": "Imagen vacía"}

        from .notebooklm_tasks import notebooklm_tasks as _tasks
        result = await _tasks.analyze_image(image_bytes, mime_type, prompt)

        # Contar como resumen/análisis en métricas del admin
        global _summaries_generated
        _summaries_generated += 1
        notebooklm_tasks.notebooklm_tasks.summaries_count += 1
        try:
            from .supabase_storage import increment_summaries_count
            increment_summaries_count()
        except Exception:
            pass

        return result
    except Exception as e:
        logging.error(f"❌ Error en /ai/analyze-image: {e}")
        return {"error": f"Error analizando imagen: {str(e)}"}


# ─── Generación de imágenes ───

@app.post("/ai/generate-image")
async def generate_image_endpoint(prompt: str = Form(...)) -> dict:
    """
    Genera una imagen a partir de un prompt usando Imagen 3 o Gemini.

    Args:
        prompt: Descripción de la imagen a generar

    Returns:
        { "image_base64": str, "mime_type": str } o { "error": str }
    """
    try:
        if not prompt or not prompt.strip():
            return {"error": "El prompt no puede estar vacío"}

        from .notebooklm_tasks import notebooklm_tasks as _tasks
        result = await _tasks.generate_image(prompt.strip())
        return result
    except Exception as e:
        logging.error(f"❌ Error en /ai/generate-image: {e}")
        return {"error": f"Error generando imagen: {str(e)}"}


@app.post("/ai/extract-text")
async def extract_text_from_file(file: UploadFile = File(...)) -> dict:
    """
    Extrae texto de archivos: TXT, PDF, DOCX, PPTX.
    
    Args:
        file: Archivo a procesar
    
    Returns:
        {
            "text": "Texto extraído",
            "word_count": número de palabras,
            "error": "null o mensaje de error"
        }
    """
    import tempfile
    from pathlib import Path
    
    try:
        # Obtener extensión
        filename = file.filename.lower()
        
        # Leer contenido
        contents = await file.read()
        
        # Procesar según extensión
        if filename.endswith('.txt'):
            text = contents.decode('utf-8')
            
        elif filename.endswith('.pdf'):
            try:
                import PyPDF2 as pypdf
                # Guardar temporalmente
                with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as tmp:
                    tmp.write(contents)
                    tmp_path = tmp.name
                
                # Extraer texto
                text = ""
                with open(tmp_path, 'rb') as pdf_file:
                    pdf_reader = pypdf.PdfReader(pdf_file)
                    for page in pdf_reader.pages:
                        text += page.extract_text() + "\n"
                
                # Limpiar
                Path(tmp_path).unlink()
                
            except Exception as e:
                return {"error": f"Error procesando PDF: {str(e)}", "text": "", "word_count": 0}
                
        elif filename.endswith('.docx'):
            try:
                from docx import Document
                # Guardar temporalmente
                with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
                    tmp.write(contents)
                    tmp_path = tmp.name
                
                # Extraer texto
                doc = Document(tmp_path)
                text = "\n".join([para.text for para in doc.paragraphs])
                
                # Limpiar
                Path(tmp_path).unlink()
                
            except Exception as e:
                return {"error": f"Error procesando DOCX: {str(e)}", "text": "", "word_count": 0}
                
        elif filename.endswith('.pptx'):
            try:
                from pptx import Presentation
                # Guardar temporalmente
                with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
                    tmp.write(contents)
                    tmp_path = tmp.name
                
                # Extraer texto
                prs = Presentation(tmp_path)
                text = ""
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
                
                # Limpiar
                Path(tmp_path).unlink()
                
            except Exception as e:
                return {"error": f"Error procesando PPTX: {str(e)}", "text": "", "word_count": 0}
        else:
            return {"error": "Formato no soportado. Usa: TXT, PDF, DOCX o PPTX", "text": "", "word_count": 0}
        
        # Contar palabras
        word_count = len(text.split())
        
        return {
            "text": text.strip(),
            "word_count": word_count,
            "error": None
        }
        
    except Exception as e:
        return {"error": f"Error procesando archivo: {str(e)}", "text": "", "word_count": 0}


@app.post("/stats", response_model=DashboardStats)
def update_stats(stats: DashboardStats) -> DashboardStats:
    """Persiste las estadísticas manuales del tablero cuando el cliente las envía."""
    save_stats(stats)
    return stats


@app.put("/api/update-streak", response_model=dict)
def update_user_streak(streak_days: int, authorization: Optional[str] = Header(None)) -> dict:
    """Actualiza las rachas del usuario actual.
    
    Parámetro: streak_days - número de días de racha
    Header: Authorization - Bearer <session_token>
    """
    if not authorization or not authorization.startswith("Bearer "):
        raise HTTPException(status_code=401, detail="Missing or invalid authorization")
    
    session_token = authorization[7:]
    session_data = validate_session_token(session_token)
    if not session_data:
        raise HTTPException(status_code=401, detail="Invalid session token")
    
    user_email = session_data.get("email")
    users = _get_cached_users()
    
    # Buscar y actualizar el usuario
    for index, user in enumerate(users):
        if user.email == user_email:
            user.streak_days = max(0, streak_days)  # No permitir valores negativos
            users[index] = user
            save_users_cached(users)
            logger.info(f"✅ Racha actualizada para {user_email}: {streak_days} días")
            return {"success": True, "streak_days": streak_days, "email": user_email}
    
    raise HTTPException(status_code=404, detail="User not found")


# ════════════════════════════════════════════════════════════════════
# FEEDBACK DE USUARIOS:
#   POST /user-feedback               → guarda calificación (1-5 estrellas)
#                                        y comentario sobre un logro. Se guarda
#                                        permanentemente y nunca se sobreescribe.
#   GET  /user-feedback/check/{id}    → verifica si el usuario ya calificó ese logro
#   GET  /admin/feedback              → lista todo el feedback (solo admin)
#   GET  /admin/feedback/stats        → estadísticas agregadas del feedback (admin)
# ════════════════════════════════════════════════════════════════════
# ==================== USER FEEDBACK ====================

@app.post("/user-feedback")
def submit_user_feedback(
    feedback: UserFeedback,
    authorization: Optional[str] = Header(None),
) -> dict:
    """Guarda el feedback (calificación + comentario) de un usuario sobre un logro.
    
    El usuario debe estar autenticado (Bearer token).
    El feedback se guarda PERMANENTEMENTE en la base de datos (nunca se borra).
    
    Parámetros (JSON body):
    - achievement_id: ID del logro calificado
    - rating: Calificación de 1-5 estrellas
    - comment: Comentario opcional (máx 200 caracteres)
    
    Header:
    - Authorization: Bearer <session_token>
    """
    # Validar autenticación
    if not authorization or not authorization.startswith("Bearer "):
        logger.warning("❌ Feedback submission without authorization")
        raise HTTPException(status_code=401, detail="Missing or invalid authorization")
    
    session_token = authorization[7:]
    session_data = validate_session_token(session_token)
    if not session_data:
        logger.warning("❌ Feedback submission with invalid session token")
        raise HTTPException(status_code=401, detail="Invalid session token")
    
    user_email = session_data.get("email")
    
    # Los datos ya vienen validados por Pydantic
    achievement_id = feedback.achievement_id
    rating_int = feedback.rating
    comment = feedback.comment or ""
    
    try:
        # Obtener nombre del usuario de la sesión
        user_name = session_data.get("display_name", "Anonymous")
        
        # Guardar feedback
        save_user_feedback(
            user_email=user_email,
            user_name=user_name,
            achievement_id=achievement_id,
            rating=rating_int,
            comment=comment
        )
        
        logger.info(f"✅ Feedback saved for {user_email}: {achievement_id} ({rating_int}/5)")
        
        return {
            "success": True,
            "message": "Gracias por tu feedback",
            "rating": rating_int,
            "achievement_id": achievement_id
        }
        
    except Exception as e:
        logger.error(f"❌ Error saving feedback: {e}")
        import traceback
        logger.error(traceback.format_exc())
        raise HTTPException(status_code=500, detail="Error saving feedback")


@app.get("/user-feedback/check/{achievement_id}")
def check_user_feedback(
    achievement_id: str,
    authorization: Optional[str] = Header(None),
) -> dict:
    """Verifica si el usuario ya ha calificado un logro específico.
    
    Retorna: {"has_rated": true/false, "achievement_id": "..."}
    """
    # Validar autenticación
    if not authorization or not authorization.startswith("Bearer "):
        raise HTTPException(status_code=401, detail="Missing or invalid authorization")
    
    session_token = authorization[7:]
    session_data = validate_session_token(session_token)
    if not session_data:
        raise HTTPException(status_code=401, detail="Invalid session token")
    
    user_email = session_data.get("email")
    
    try:
        if check_user_feedback_exists is None:
            raise HTTPException(status_code=503, detail="Feedback service not available")
        
        has_rated = check_user_feedback_exists(user_email, achievement_id)
        logger.info(f"📋 Checked feedback for {user_email} + {achievement_id}: {has_rated}")
        
        return {
            "has_rated": has_rated,
            "achievement_id": achievement_id,
            "user_email": user_email
        }
        
    except Exception as e:
        logger.error(f"❌ Error checking feedback: {e}")
        raise HTTPException(status_code=500, detail="Error checking feedback")


@app.get("/admin/feedback")
def get_admin_feedback(
    authorization: Optional[str] = Header(None),
) -> list:
    """Retorna TODOS los feedback de usuarios para el panel administrativo.
    
    Los datos NUNCA se borran - son datos valiosos para reportes y análisis.
    
    Requiere autenticación de admin.
    """
    # Validar autenticación
    if not authorization or not authorization.startswith("Bearer "):
        logger.warning("❌ Admin feedback request without authorization")
        raise HTTPException(status_code=401, detail="Missing or invalid authorization")
    
    session_token = authorization[7:]
    session_data = validate_session_token(session_token)
    if not session_data:
        logger.warning("❌ Admin feedback request with invalid token")
        raise HTTPException(status_code=401, detail="Invalid session token")
    
    user_email = session_data.get("email")
    
    # TODO: Implementar verificación de rol admin (por ahora permite a cualquier usuario autenticado)
    # En el futuro, verifica si el usuario tiene rol "admin"
    
    try:
        if load_all_user_feedback is None:
            raise HTTPException(status_code=503, detail="Feedback service not available")
        
        feedback_list = load_all_user_feedback()
        logger.info(f"✅ Admin retrieved {len(feedback_list)} feedback records")
        
        return feedback_list
        
    except Exception as e:
        logger.error(f"❌ Error loading admin feedback: {e}")
        import traceback
        logger.error(traceback.format_exc())
        raise HTTPException(status_code=500, detail="Error loading feedback data")


@app.get("/admin/feedback/stats")
def get_admin_feedback_stats(
    authorization: Optional[str] = Header(None),
) -> dict:
    """Retorna estadísticas agregadas del feedback (promedio, distribución, etc).
    """
    # Validar autenticación
    if not authorization or not authorization.startswith("Bearer "):
        raise HTTPException(status_code=401, detail="Missing or invalid authorization")
    
    session_token = authorization[7:]
    session_data = validate_session_token(session_token)
    if not session_data:
        raise HTTPException(status_code=401, detail="Invalid session token")
    
    try:
        if get_feedback_stats is None:
            raise HTTPException(status_code=503, detail="Feedback service not available")
        
        stats = get_feedback_stats()
        logger.info(f"✅ Feedback stats retrieved: {stats}")
        
        return stats
        
    except Exception as e:
        logger.error(f"❌ Error calculating feedback stats: {e}")
        raise HTTPException(status_code=500, detail="Error calculating feedback stats")


@app.delete("/admin/feedback/{feedback_id}")
def delete_admin_feedback(
    feedback_id: int,
    authorization: Optional[str] = Header(None),
) -> dict:
    """Elimina una calificación de usuario (solo admin)."""
    if not authorization or not authorization.startswith("Bearer "):
        raise HTTPException(status_code=401, detail="Missing or invalid authorization")

    session_token = authorization[7:]
    session_data = validate_session_token(session_token)
    if not session_data:
        raise HTTPException(status_code=401, detail="Invalid session token")

    admin_email = os.getenv("ADMIN_EMAIL", "appscolyax@gmail.com")
    if session_data.get("email") != admin_email:
        raise HTTPException(status_code=403, detail="Forbidden - Admin access required")

    try:
        from .user_feedback import delete_user_feedback
        deleted = delete_user_feedback(feedback_id)
        if not deleted:
            raise HTTPException(status_code=404, detail="Calificación no encontrada")
        return {"success": True, "deleted_id": feedback_id}
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"❌ Error eliminando feedback {feedback_id}: {e}")
        raise HTTPException(status_code=500, detail="Error eliminando calificación")


# ============================================================================
# ADMIN ENDPOINTS - Gestión de usuarios y estadísticas
# ============================================================================

@app.get("/admin/metrics")
def get_admin_metrics(authorization: Optional[str] = Header(None)) -> dict:
    """Retorna métricas del sistema para el dashboard administrativo."""
    # Validar autenticación
    if not authorization or not authorization.startswith("Bearer "):
        raise HTTPException(status_code=401, detail="Missing or invalid authorization")
    
    session_token = authorization[7:]
    session_data = validate_session_token(session_token)
    if not session_data:
        raise HTTPException(status_code=401, detail="Invalid session token")
    
    # Verificar que es admin (email específico)
    admin_email = os.getenv("ADMIN_EMAIL", "appscolyax@gmail.com")
    if session_data.get("email") != admin_email:
        logger.warning(f"❌ Non-admin user attempted access: {session_data.get('email')}")
        raise HTTPException(status_code=403, detail="Forbidden - Admin access required")
    
    from .supabase_storage import load_users, load_tasks

    # ── Usuarios totales ──
    total_users = 0
    try:
        total_users = len(load_users())
    except Exception as e:
        logger.warning(f"⚠️ No se pudo contar usuarios: {e}")

    # ── Tareas completadas ──
    tasks_completed = 0
    try:
        all_tasks = load_tasks()
        tasks_completed = sum(1 for t in all_tasks if t.status == TaskStatus.COMPLETED)
    except Exception as e:
        logger.warning(f"⚠️ No se pudo contar tareas completadas: {e}")

    # ── Resúmenes generados (tabla summaries vía Supabase directo) ──
    summaries_generated = _summaries_generated  # contador en memoria como respaldo
    try:
        sb = get_supabase()
        if sb:
            summaries_data = sb.table("summaries").select("id").execute().data or []
            summaries_generated = len(summaries_data)
    except Exception as e:
        logger.warning(f"⚠️ No se pudo contar summaries: {e}")

    # ── Streaks activos (user_stats) ──
    active_streaks = 0
    try:
        sb = get_supabase()
        if sb:
            user_stats = sb.table("user_stats").select("streak_days").execute().data or []
            active_streaks = sum(1 for s in user_stats if (s.get("streak_days") or 0) > 0)
    except Exception as e:
        logger.warning(f"⚠️ No se pudo contar streaks: {e}")

    retention_rate = 60  # aproximación fija hasta implementar cálculo real

    metrics = {
        "total_users": total_users,
        "active_users_30d": total_users,
        "summaries_generated": summaries_generated,
        "active_streaks": active_streaks,
        "retention_rate": retention_rate,
        "tasks_completed": tasks_completed,
    }

    logger.info(f"✅ Admin metrics: users={total_users}, tasks_completed={tasks_completed}, summaries={summaries_generated}, streaks={active_streaks}")
    return metrics


@app.get("/admin/iris-stats")
def get_iris_stats(authorization: Optional[str] = Header(None)) -> dict:
    """Retorna estadísticas de uso de Iris IA para el panel administrativo."""
    if not authorization or not authorization.startswith("Bearer "):
        raise HTTPException(status_code=401, detail="Missing or invalid authorization")

    session_token = authorization[7:]
    session_data = validate_session_token(session_token)
    if not session_data:
        raise HTTPException(status_code=401, detail="Invalid session token")

    admin_email = os.getenv("ADMIN_EMAIL", "appscolyax@gmail.com")
    if session_data.get("email") != admin_email:
        raise HTTPException(status_code=403, detail="Forbidden - Admin access required")

    # Resúmenes generados — leer de Supabase (persiste entre deploys)
    summaries_count = _summaries_generated  # fallback en memoria
    try:
        from .supabase_storage import get_summaries_count
        db_count = get_summaries_count()
        # Usar el mayor entre DB y memoria (DB es más confiable)
        summaries_count = max(db_count, _summaries_generated)
    except Exception as e:
        logger.warning(f"⚠️ No se pudo leer summaries de Supabase: {e}")

    # Sesiones de enfoque — contar todas las del sistema
    focus_sessions_total = 0
    try:
        from .supabase_storage import get_focus_sessions_count
        focus_sessions_total = get_focus_sessions_count()
    except Exception as e:
        logger.warning(f"⚠️ No se pudo contar focus_sessions: {e}")

    # Total de tareas en el sistema
    tasks_total = 0
    try:
        from .supabase_storage import load_tasks
        tasks_total = len(load_tasks())
    except Exception as e:
        logger.warning(f"⚠️ No se pudo contar tasks: {e}")

    # Usuarios activos en tiempo real (caché en memoria + DB)
    realtime_users = 0
    try:
        from .supabase_storage import count_realtime_users
        realtime_users = count_realtime_users()
    except Exception as e:
        logger.warning(f"⚠️ No se pudo contar realtime_users: {e}")

    # Estadísticas en memoria de la IA + acumulado en Supabase
    ai_requests = 0
    cache_hits = 0
    requests_remaining = 20
    try:
        from .notebooklm_tasks import notebooklm_tasks as _nt
        ai_stats = _nt.get_stats()
        ai_requests = ai_stats.get("requests_today", 0)
        requests_remaining = ai_stats.get("requests_remaining", 20)
    except Exception as e:
        logger.error(f"❌ Error obteniendo stats de IA: {e}")
    
    # cache_hits acumulado desde Supabase (persiste entre deploys)
    try:
        from .supabase_storage import get_cache_hits
        cache_hits = get_cache_hits()
        logger.info(f"✅ Cache hits obtenidos: {cache_hits}")
    except Exception as e:
        logger.error(f"❌ Error obteniendo cache_hits: {e}")

    logger.info(
        f"📊 Iris stats: summaries={summaries_count}, tasks={tasks_total}, "
        f"focus={focus_sessions_total}, realtime={realtime_users}, "
        f"ai_requests={ai_requests}, cache_hits={cache_hits}"
    )

    return {
        "summaries_count": summaries_count,
        "focus_sessions_total": focus_sessions_total,
        "tasks_total": tasks_total,
        "realtime_users": realtime_users,
        "ai_requests_session": ai_requests,
        "cache_hits": cache_hits,
        "requests_remaining": requests_remaining,
    }


# ─── Testing endpoint para verificar cache hits ───
@app.post("/admin/test-cache-hit")
async def test_cache_hit(
    question: str = Query(...),
    authorization: Optional[str] = Header(None)
) -> dict:
    """Endpoint para testing: verifica que el cache funcione y cuente hits.
    
    Uso:
    1. Primera llamada: POST /admin/test-cache-hit?question=¿Cuál%20es%20el%20ciclo%20del%20agua?
       → Genera respuesta de IA, la guarda en caché
    2. Segunda llamada: misma pregunta
       → Devuelve del caché, incrementa contador
    
    Respuesta:
    {
        "question": "Tu pregunta",
        "answer": "Respuesta de Iris",
        "cache_hits": 5  # Contador actual
    }
    """
    if not authorization or not authorization.startswith("Bearer "):
        raise HTTPException(status_code=401, detail="Admin access required")
    
    session_token = authorization[7:]
    session_data = validate_session_token(session_token)
    if not session_data:
        raise HTTPException(status_code=401, detail="Invalid session token")
    
    admin_email = os.getenv("ADMIN_EMAIL", "appscolyax@gmail.com")
    if session_data.get("email") != admin_email:
        raise HTTPException(status_code=403, detail="Forbidden - Admin access required")
    
    if not question or not question.strip():
        raise HTTPException(status_code=400, detail="Question cannot be empty")
    
    try:
        # Simular Q&A para testing del caché
        document = "Scolyax es una plataforma académica para estudiantes con ADHD. El caché es una técnica de optimización."
        result = await notebooklm_tasks.answer_question(document, question)
        
        # Obtener contador después
        from .supabase_storage import get_cache_hits
        current_cache_hits = get_cache_hits()
        
        logger.info(f"✅ Test cache: question='{question}' | answer_len={len(result.get('answer', ''))} | cache_hits={current_cache_hits}")
        
        return {
            "question": question,
            "answer": result.get("answer", ""),
            "cache_hits": current_cache_hits
        }
    except Exception as e:
        logger.error(f"❌ Error en test-cache-hit: {e}", exc_info=True)
        raise HTTPException(status_code=500, detail=str(e))


@app.get("/api/admin/users")
def list_admin_users(
    page: int = 1,
    limit: int = 20,
    authorization: Optional[str] = Header(None)
) -> dict:
    """Retorna lista de usuarios para el panel administrativo."""
    # Validar autenticación
    if not authorization or not authorization.startswith("Bearer "):
        raise HTTPException(status_code=401, detail="Missing or invalid authorization")
    
    session_token = authorization[7:]
    session_data = validate_session_token(session_token)
    if not session_data:
        raise HTTPException(status_code=401, detail="Invalid session token")
    
    # Verificar que es admin
    admin_email = os.getenv("ADMIN_EMAIL", "appscolyax@gmail.com")
    if session_data.get("email") != admin_email:
        logger.warning(f"❌ Non-admin access attempt from: {session_data.get('email')}")
        raise HTTPException(status_code=403, detail="Forbidden - Admin access required")
    
    try:
        sb = get_supabase()
        
        # Obtener usuarios usando Supabase
        all_users = sb.table("users").select("*").execute().data or []
        
        # Obtener stats
        all_stats = sb.table("user_stats").select("*").execute().data or []
        stats_by_user = {s.get("user_id"): s for s in all_stats}
        
        # Ordenar por fecha creación descendente
        all_users.sort(key=lambda u: u.get("created_at", ""), reverse=True)
        
        # Pagination
        offset = (page - 1) * limit
        paginated_users = all_users[offset:offset + limit]
        
        users = []
        for user in paginated_users:
            user_id = str(user.get("id", ""))
            stats = stats_by_user.get(user_id, {})
            
            users.append({
                "id": user_id,
                "email": user.get("email", ""),
                "display_name": user.get("display_name") or "Sin nombre",
                "provider": user.get("provider", "email"),
                "created_at": user.get("created_at"),
                "streak_days": stats.get("streak_days", 0) or 0,
                "total_xp": stats.get("total_xp", 0) or 0,
                "level": stats.get("level", 1) or 1
            })
        
        logger.info(f"✅ Admin retrieved {len(users)} users (page {page})")
        
        return {
            "users": users,
            "total": len(all_users),
            "page": page,
            "limit": limit,
            "pages": (len(all_users) + limit - 1) // limit
        }
        
    except Exception as e:
        logger.error(f"❌ Error listing users: {e}")
        import traceback
        logger.error(traceback.format_exc())
        raise HTTPException(status_code=500, detail="Error listing users")


@app.delete("/api/admin/users/{user_id}")
def delete_admin_user(
    user_id: str,
    authorization: Optional[str] = Header(None)
) -> dict:
    """Elimina un usuario (solo admin)."""
    # Validar autenticación
    if not authorization or not authorization.startswith("Bearer "):
        raise HTTPException(status_code=401, detail="Missing or invalid authorization")
    
    session_token = authorization[7:]
    session_data = validate_session_token(session_token)
    if not session_data:
        raise HTTPException(status_code=401, detail="Invalid session token")
    
    # Verificar que es admin
    admin_email = os.getenv("ADMIN_EMAIL", "appscolyax@gmail.com")
    if session_data.get("email") != admin_email:
        raise HTTPException(status_code=403, detail="Forbidden - Admin access required")
    
    # Prevenir auto-eliminación
    if str(session_data.get("user_id")) == str(user_id):
        raise HTTPException(status_code=400, detail="No puedes eliminarte a ti mismo")
    
    try:
        sb = get_supabase()
        
        # Obtener usuario para confirmar que existe
        user_response = sb.table("users").select("id, email").eq("id", user_id).execute()
        if not user_response.data:
            raise HTTPException(status_code=404, detail="Usuario no encontrado")
        
        user_email = user_response.data[0].get("email", "Unknown")
        
        # Eliminar todos los datos del usuario en cascada
        logger.info(f"🗑️ Iniciando eliminación en cascada para {user_email}")
        
        try:
            sb.table("tasks").delete().eq("user_email", user_email).execute()
            logger.info(f"✅ Eliminadas tareas de {user_email}")
        except Exception as e:
            logger.warning(f"⚠️ Error eliminando tareas: {e}")
        
        try:
            sb.table("reminders").delete().eq("user_email", user_email).execute()
            logger.info(f"✅ Eliminados recordatorios de {user_email}")
        except Exception as e:
            logger.warning(f"⚠️ Error eliminando recordatorios: {e}")
        
        try:
            sb.table("focus_sessions").delete().eq("user_email", user_email).execute()
            logger.info(f"✅ Eliminadas sesiones de enfoque de {user_email}")
        except Exception as e:
            logger.warning(f"⚠️ Error eliminando sesiones de enfoque: {e}")
        
        try:
            sb.table("user_sessions").delete().eq("email", user_email).execute()
            logger.info(f"✅ Eliminadas sesiones del usuario {user_email}")
        except Exception as e:
            logger.warning(f"⚠️ Error eliminando sesiones: {e}")
        
        try:
            sb.table("user_stats").delete().eq("user_email", user_email).execute()
            logger.info(f"✅ Eliminadas estadísticas de {user_email}")
        except Exception as e:
            logger.warning(f"⚠️ Error eliminando estadísticas: {e}")
        
        try:
            sb.table("user_achievements").delete().eq("user_email", user_email).execute()
            logger.info(f"✅ Eliminados logros de {user_email}")
        except Exception as e:
            logger.warning(f"⚠️ Error eliminando logros: {e}")
        
        try:
            sb.table("user_feedback").delete().eq("user_email", user_email).execute()
            logger.info(f"✅ Eliminado feedback de {user_email}")
        except Exception as e:
            logger.warning(f"⚠️ Error eliminando feedback: {e}")
        
        try:
            sb.table("summaries").delete().eq("user_id", user_id).execute()
            logger.info(f"✅ Eliminados resúmenes de {user_email}")
        except Exception as e:
            logger.warning(f"⚠️ Error eliminando resúmenes: {e}")
        
        # Finalmente, eliminar el usuario
        sb.table("users").delete().eq("id", user_id).execute()
        logger.info(f"✅ Eliminado usuario {user_email} completamente")
        
        return {
            "success": True,
            "message": f"Usuario {user_email} eliminado correctamente con todos sus datos",
            "deleted_user_id": user_id
        }
        
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"❌ Error deleting user: {e}")
        import traceback
        logger.error(traceback.format_exc())
        raise HTTPException(status_code=500, detail="Error eliminating user")


@app.put("/api/admin/users/{user_id}")
def update_admin_user(
    user_id: str,
    update_data: dict,
    authorization: Optional[str] = Header(None)
) -> dict:
    """Actualiza datos de un usuario (solo admin)."""
    # Validar autenticación
    if not authorization or not authorization.startswith("Bearer "):
        raise HTTPException(status_code=401, detail="Missing or invalid authorization")
    
    session_token = authorization[7:]
    session_data = validate_session_token(session_token)
    if not session_data:
        raise HTTPException(status_code=401, detail="Invalid session token")
    
    # Verificar que es admin
    admin_email = os.getenv("ADMIN_EMAIL", "appscolyax@gmail.com")
    if session_data.get("email") != admin_email:
        raise HTTPException(status_code=403, detail="Forbidden - Admin access required")
    
    try:
        sb = get_supabase()
        
        # Actualizar usuario
        if "display_name" in update_data:
            sb.table("users").update({"display_name": update_data["display_name"]}).eq("id", user_id).execute()
        
        # Actualizar stats
        stats_update = {}
        if "streak_days" in update_data:
            stats_update["streak_days"] = update_data["streak_days"]
        if "total_xp" in update_data:
            stats_update["total_xp"] = update_data["total_xp"]
        if "level" in update_data:
            stats_update["level"] = update_data["level"]
        
        if stats_update:
            sb.table("user_stats").update(stats_update).eq("user_id", user_id).execute()
        
        logger.info(f"✅ Admin updated user: {user_id}")
        
        return {
            "success": True,
            "message": "Usuario actualizado correctamente",
            "user_id": user_id
        }
        
    except Exception as e:
        logger.error(f"❌ Error updating user: {e}")
        import traceback
        logger.error(traceback.format_exc())
        raise HTTPException(status_code=500, detail="Error updating user")


# ════════════════════════════════════════════════════════════════════
# PANEL DE ADMINISTRACIÓN:
#   GET    /admin/metrics            → métricas globales del sistema (usuarios,
#                                      tareas, sesiones) — solo admin
#   GET    /api/admin/users          → lista paginada de usuarios con sus stats
#   DELETE /api/admin/users/{id}    → elimina un usuario y TODOS sus datos
#                                      (tareas, recordatorios, logros, feedback)
#   PUT    /api/admin/users/{id}    → actualiza nombre, racha, XP o nivel
#
# Todos los endpoints admin validan que el Bearer token pertenezca
# al email configurado en la variable de entorno ADMIN_EMAIL.
# ════════════════════════════════════════════════════════════════════
# ============================================================================
# TTS (Text-to-Speech) - Voces realistas con Google Cloud
# ============================================================================

# ════════════════════════════════════════════════════════════════════
# TTS (TEXT-TO-SPEECH) — Voz de Iris:
#   POST /tts          → síntesis completa con Google Cloud Neural TTS.
#                        Devuelve audio MP3. Voces neurales femeninas en es-ES.
#   GET  /tts/stream   → síntesis en streaming via edge-tts (Microsoft Neural).
#                        El audio empieza a reproducirse antes de terminar.
#   GET  /tts/voices   → lista las voces disponibles para un idioma.
#
# Prioridad: Google Cloud TTS (Neural) → edge-tts (Microsoft) → gTTS (fallback).
# Límite libre: 1 millón de caracteres/mes en Google Cloud TTS.
# ════════════════════════════════════════════════════════════════════
@app.post("/tts")
async def synthesize_speech(
    text: str = Form(...),
    language: str = Form(default="es-ES"),
    voice_name: str = Form(default="es-ES-Neural2-A"),
    speaking_rate: float = Form(default=0.95),
    pitch: float = Form(default=-2.0)
) -> Response:
    """
    Sintetiza texto a audio usando Google Cloud Text-to-Speech (voces realistas).
    
    **Características:**
    - Voces IA naturales y humanas de alta calidad
    - Voces Neural de Google (más naturales que WaveNet)
    - Completamente gratis (1M caracteres/mes)
    - Voz femenina realista en español (Iris)
    - Audio MP3 de alta calidad
    - Control de velocidad y tono para mayor naturalidad
    
    **Parámetros:**
    - text: Texto a convertir a audio (max 5000 caracteres)
    - language: Código de idioma (default: es-ES para español)
    - voice_name: Nombre de la voz (default: es-ES-Neural2-A, voz femenina natural)
    - speaking_rate: Velocidad de habla (0.25-4.0, default: 0.95 para mayor claridad)
    - pitch: Tono de voz (-20.0 a 20.0, default: -2.0 para voz más cálida)
    
    **Voces disponibles en español (femeninas):**
    - es-ES-Neural2-A: Mujer, natural, clara (RECOMENDADO - Iris)
    - es-ES-Neural2-C: Mujer, natural, profesional
    - es-ES-Neural2-E: Mujer, natural, cálida
    - es-ES-Wavenet-C: Mujer, wavenet (muy realista)
    
    **Respuesta:**
    - Audio MP3 en el body con Content-Type: audio/mpeg
    
    **Ejemplo:**
    ```bash
    curl -X POST "http://localhost:8000/tts" \\
      -F "text=Hola, soy Iris, tu asistente académica inteligente" \\
      -F "language=es-ES" \\
      -F "voice_name=es-ES-Neural2-A" \\
      -F "speaking_rate=0.95" \\
      -F "pitch=-2.0" \\
      --output audio.mp3
    ```
    """
    
    # Validar entrada
    if not text or not text.strip():
        raise HTTPException(status_code=400, detail="Text cannot be empty")
    
    if len(text) > 5000:
        logger.warning(f"⚠️  TTS text too long ({len(text)} chars), truncating to 5000")
        text = text[:5000]
    
    try:
        # Obtener audio sintetizado con parámetros de voz personalizados
        audio_bytes = await tts_provider.get_tts_audio(
            text=text.strip(),
            language=language,
            voice_name=voice_name,
            speaking_rate=speaking_rate,
            pitch=pitch,
            use_google=True
        )
        
        if not audio_bytes:
            logger.warning("⚠️  Google Cloud TTS returned empty audio")
            raise HTTPException(
                status_code=503,
                detail="Text-to-Speech service not available. Make sure Google Cloud credentials are configured."
            )
        
        logger.info(f"✅ TTS synthesized: {len(text)} chars → {len(audio_bytes)} bytes of audio")
        
        # Retorna audio MP3
        return Response(
            content=audio_bytes,
            media_type="audio/mpeg",
            headers={"Content-Disposition": "attachment; filename=audio.mp3"}
        )
        
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"❌ TTS synthesis error: {e}")
        import traceback
        logger.error(traceback.format_exc())
        raise HTTPException(status_code=500, detail="Error synthesizing speech")


@app.get("/tts/stream")
async def tts_stream(
    text: str,
    voice: str = "es-CO-SalomeNeural",
    rate: str = "+0%",
    pitch: str = "+0Hz"
):
    """
    TTS con streaming via edge-tts. El audio empieza a reproducirse inmediatamente
    mientras se sigue generando. Voces neuronales humanas de Microsoft.
    
    Voces recomendadas (femeninas, humanas):
    - es-CO-SalomeNeural (colombiana, muy natural)
    - es-MX-DaliaNeural (mexicana, cálida)
    - es-ES-ElviraNeural (española, clara)
    """
    if not text or not text.strip():
        raise HTTPException(status_code=400, detail="Text cannot be empty")
    
    cleaned = text.strip()[:5000]
    
    try:
        import edge_tts
    except ImportError:
        raise HTTPException(status_code=503, detail="edge-tts not installed")
    
    async def audio_stream():
        communicate = edge_tts.Communicate(cleaned, voice=voice, rate=rate, pitch=pitch)
        async for chunk in communicate.stream():
            if chunk["type"] == "audio":
                yield chunk["data"]
    
    logger.info(f"🔊 TTS stream: {len(cleaned)} chars, voice={voice}")
    return StreamingResponse(
        audio_stream(),
        media_type="audio/mpeg",
        headers={
            "Cache-Control": "no-cache",
            "Transfer-Encoding": "chunked"
        }
    )


@app.get("/tts/voices")
async def get_available_voices(language: str = "es-ES") -> dict:
    """
    Retorna las voces disponibles para un idioma.
    
    **Parámetros:**
    - language: Código de idioma (default: es-ES)
    
    **Respuesta:**
    ```json
    {
        "language": "es-ES",
        "available": true,
        "voices": ["es-ES-Neural2-a", "es-ES-Neural2-c", "es-ES-Standard-a", ...],
        "recommended_male": "es-ES-Neural2-c"
    }
    ```
    """
    try:
        voices = await tts_provider.google_tts.get_available_voices(language)
        
        # Identifica voces masculinas recomendadas
        recommended_male = "es-ES-Neural2-c" if language == "es-ES" else None
        
        return {
            "language": language,
            "available": tts_provider.google_tts.is_available(),
            "voices": voices,
            "recommended_male": recommended_male,
            "info": "Use recommended_male for realistic male voice"
        }
        
    except Exception as e:
        logger.error(f"❌ Error listing TTS voices: {e}")
        return {
            "language": language,
            "available": False,
            "voices": [],
            "recommended_male": None,
            "error": str(e)
        }


# ════════════════════════════════════════════════════════════════════
# SISTEMA IA DE GESTIÓN DEL TIEMPO (ai_time_management):
#   POST /ai/test/analyze     → analiza respuestas del test cognitivo y recomienda
#                                una de 3 herramientas de gestión del tiempo.
#   POST /ai/session/create   → crea una sesión de estudio con la herramienta
#                                recomendada y genera checkpoints automáticos.
#
# Este módulo detecta el estilo de aprendizaje del estudiante (secuencial,
# adaptable, visual...) y asigna el método más adecuado:
#   - Pomodoro: para perfiles con alta distracción
#   - Time blocking: para perfiles con mayor autocontrol
#   - Adaptive learning: para perfiles flexibles
# ════════════════════════════════════════════════════════════════════
# ==================== SISTEMA IA DE GESTIÓN DE TIEMPO ====================

@app.post("/ai/test/analyze", response_model=AIToolRecommendation)
async def analyze_test_and_recommend(
    test_data: dict,
    authorization: Optional[str] = Header(None)
):
    """
    Analiza el test del usuario y recomienda 1 de 3 herramientas de gestión de tiempo.
    
    **Parámetros:**
    - test_data: Objeto JSON con respuestas del test
    
    **Flujo:**
    1. Valida sesión del usuario
    2. Analiza respuestas del test
    3. Detecta estilo de aprendizaje
    4. Recomienda herramienta óptima
    5. Retorna recomendación con confianza
    
    **Respuesta:**
    ```json
    {
        "tool_type": "adaptive_learning",
        "confidence": 0.87,
        "reasoning": "Tu estilo es adaptable...",
        "estimated_session_length": 50,
        "user_email": "user@example.com"
    }
    ```
    """
    try:
        # Validar token
        session = validate_session_token(authorization)
        if not session:
            raise HTTPException(status_code=401, detail="No session")
        
        user_email = session.user_email
        
        # Analizar test
        test_result = ai_time_management.analyze_test_results(test_data, user_email)
        
        # Recomendar herramienta
        recommendation = ai_time_management.recommend_ai_tool(test_result, user_email)
        
        # Guardar recomendación (opcional, en storage)
        # TODO: Guardar en tabla recommendations en Supabase
        
        logger.info(f"✅ Test analizado para {user_email}: {recommendation.tool_type.value}")
        
        return recommendation
        
    except Exception as e:
        logger.error(f"❌ Error anaizando test: {e}")
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/ai/session/create", response_model=AIStudySession)
async def create_ai_study_session(
    task_id: int,
    tool_type: AITimeManagementTool,
    authorization: Optional[str] = Header(None)
):
    """
    Crea una sesión de estudio con herramienta de IA y genera checkpoints.
    
    **Parámetros:**
    - task_id: ID de la tarea
    - tool_type: Una de: adaptive_learning, deep_focus, goal_tracking
    
    **Retorna:**
    Sesión con checkpoints pre-generados y programados.
    """
    try:
        # Validar sesión
        session = validate_session_token(authorization)
        if not session:
            raise HTTPException(status_code=401, detail="No session")
        
        user_email = session.user_email
        
        # Cargar tarea
        tasks = load_tasks(user_email)
        task = next((t for t in tasks if t.id == task_id), None)
        
        if not task:
            raise HTTPException(status_code=404, detail="Task not found")
        
        # Crear sesión
        session_id = str(uuid.uuid4())
        
        # Determinar duración estimada
        estimated_length = 50  # Default
        if tool_type == AITimeManagementTool.DEEP_FOCUS:
            estimated_length = 60
        elif tool_type == AITimeManagementTool.GOAL_TRACKING:
            estimated_length = task.estimated_pomodoros * 25 if task.estimated_pomodoros > 0 else 60
        
        # Generar checkpoints
        checkpoints = ai_time_management.generate_checkpoints(
            session_id=session_id,
            tool_type=tool_type,
            task_title=task.title,
            estimated_session_length=estimated_length
        )
        
        # Crear objeto sesión
        study_session = AIStudySession(
            id=session_id,
            user_email=user_email,
            task_id=task_id,
            task_title=task.title,
            tool_type=tool_type,
            checkpoints=checkpoints,
            started_at=datetime.now(timezone.utc),
            total_time_minutes=estimated_length
        )
        
        # TODO: Guardar en Supabase tabla ai_study_sessions
        
        logger.info(f"✅ Sesión IA creada: {session_id} con {len(checkpoints)} checkpoints")
        
        return study_session
        
    except Exception as e:
        logger.error(f"❌ Error creando sesión IA: {e}")
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/ai/checkpoint/submit", response_model=CheckpointVerificationResponse)
async def submit_checkpoint_evidence(
    request: CheckpointVerificationRequest,
    authorization: Optional[str] = Header(None)
):
    """
    Usuario envía evidencia de checkpoint (foto + descripción).
    Iris (IA) verifica y retorna feedback.
    
    **Parámetros:**
    - checkpoint_id: ID del checkpoint
    - session_id: ID de la sesión
    - user_description: Descripción de lo completado
    - photo_base64: Foto en base64 (opcional)
    
    **Retorna:**
    Análisis con:
    - verified: bool (aprobado o no)
    - confidence: 0-1 (confianza del análisis)
    - ai_feedback: Retroalimentación de Iris
    - suggestions: Tips para próximo checkpoint
    """
    try:
        # Validar sesión
        session = validate_session_token(authorization)
        if not session:
            raise HTTPException(status_code=401, detail="No session")
        
        user_email = session.user_email
        
        # Verificar que el email coincide
        if request.user_email != user_email:
            raise HTTPException(status_code=403, detail="Unauthorized checkpoint")
        
        # Verificar checkpoint con IA
        verification = ai_time_management.verify_checkpoint_with_ai(request)
        
        # TODO: 
        # 1. Guardar foto en S3
        # 2. Actualizar checkpoint en BD
        # 3. Enviar email al usuario
        # 4. Si aprobado, desbloquear siguiente checkpoint
        
        logger.info(f"✅ Checkpoint verificado: {verification.checkpoint_id}, verified={verification.verified}")
        
        return verification
        
    except Exception as e:
        logger.error(f"❌ Error verificando checkpoint: {e}")
        raise HTTPException(status_code=500, detail=str(e))


@app.get("/ai/recommendations")
async def get_ai_recommendations(
    authorization: Optional[str] = Header(None)
):
    """
    Obtiene las recomendaciones de herramientas post-test para el usuario actual.
    
    **Retorna:**
    ```json
    {
        "recommended_tools": [
            {
                "tool_type": "adaptive_learning",
                "confidence": 0.87,
                "reasoning": "...",
                "status": "active" | "pending" | "completed"
            }
        ],
        "test_date": "2026-02-17T10:00:00Z",
        "learning_profile": { ... }
    }
    ```
    """
    try:
        session = validate_session_token(authorization)
        if not session:
            raise HTTPException(status_code=401, detail="No session")
        
        # TODO: Cargar recomendaciones de la BD (tabla ai_recommendations)
        
        return {
            "recommended_tools": [],
            "test_date": None,
            "message": "No recommendations yet. Complete the test first."
        }
        
    except Exception as e:
        logger.error(f"❌ Error obteniendo recomendaciones: {e}")
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/ai/focus-alert")
async def focus_alert(
    request: dict,
    authorization: Optional[str] = Header(None)
):
    """
    Anti-cheat: Registra cuando el usuario cambia de pestaña durante una sesión activa.
    Opcionalmente envía una alerta por email si hay múltiples cambios.
    
    **Parámetros (JSON body):**
    - session_id: ID de la sesión de estudio
    - event_type: 'tab_switch' | 'window_blur'
    - timestamp: ISO timestamp del evento
    - warning_count: Número total de advertencias acumuladas
    - task_title: Título de la tarea activa
    
    **Retorna:** { "status": "logged", "email_sent": bool }
    """
    try:
        session = validate_session_token(authorization)
        if not session:
            raise HTTPException(status_code=401, detail="No session")
        
        user_email = session.user_email
        warning_count = request.get("warning_count", 0)
        task_title = request.get("task_title", "Sin título")
        event_type = request.get("event_type", "tab_switch")
        
        logger.warning(
            f"⚠️ Focus alert [{event_type}] - User: {user_email}, "
            f"Task: {task_title}, Warnings: {warning_count}"
        )
        
        email_sent = False
        
        # Send email alert on 3rd+ warning
        if warning_count >= 3:
            try:
                from email.message import EmailMessage as EM
                from .mailer import _send_message_or_outbox
                msg = EM()
                msg["Subject"] = "⚠️ Scolyax - Alerta de concentración"
                msg["To"] = user_email
                msg["From"] = os.getenv("SCOLYAX_EMAIL_FROM", "no-reply@scolyax.local")
                msg.set_content(
                    f"Hola,\n\n"
                    f"Se detectaron {warning_count} cambios de pestaña durante tu sesión "
                    f'de estudio en la tarea "{task_title}".\n\n'
                    f"Mantener el enfoque es clave para tu productividad. "
                    f"Intenta cerrar pestañas innecesarias y silenciar notificaciones.\n\n"
                    f"— Iris, tu asistente Scolyax 🎓"
                )
                _send_message_or_outbox(msg, prefix="focus_alert")
                email_sent = True
                logger.info(f"📧 Focus alert email sent/queued for {user_email}")
            except Exception as mail_err:
                logger.error(f"❌ Error sending focus alert email: {mail_err}")
        
        return {"status": "logged", "email_sent": email_sent, "warnings": warning_count}
        
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"❌ Error en focus alert: {e}")
        raise HTTPException(status_code=500, detail=str(e))


# ==================== ENERGY JOURNAL ====================


@app.post("/energy-journal")
async def create_energy_entry(
    entry: EnergyEntryCreate,
    authorization: str = Header(None),
):
    """Registra una entrada de energía tras completar una sesión de estudio."""
    try:
        session = _require_session(authorization)
        user_email = session.email

        from .supabase_storage import save_energy_entry

        energy_entry = EnergyEntry(
            user_email=user_email,
            energy_level=entry.energy_level,
            mood=entry.mood,
            notes=entry.notes,
            session_type=entry.session_type,
            session_duration_minutes=entry.session_duration_minutes,
        )

        saved = save_energy_entry(energy_entry, user_email)
        if saved:
            return {"status": "ok", "entry": saved.model_dump()}
        return {"status": "ok", "entry": energy_entry.model_dump()}

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"❌ Error saving energy entry: {e}")
        raise HTTPException(status_code=500, detail=str(e))


@app.get("/energy-journal")
async def get_energy_entries(
    limit: int = 30,
    authorization: str = Header(None),
):
    """Devuelve las últimas entradas de energía del usuario."""
    try:
        session = _require_session(authorization)
        user_email = session.email

        from .supabase_storage import load_energy_entries

        entries = load_energy_entries(user_email, limit=limit)
        return {
            "entries": [e.model_dump() for e in entries],
            "count": len(entries),
        }

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"❌ Error loading energy entries: {e}")
        raise HTTPException(status_code=500, detail=str(e))


# ==================== CRISIS MODE ====================


@app.post("/crisis-mode/start")
async def start_crisis_session(
    payload: CrisisSessionCreate,
    authorization: str = Header(None),
):
    """Registra el inicio de una sesión del modo crisis."""
    try:
        session = _require_session(authorization)
        user_email = session.email

        from .supabase_storage import save_crisis_session

        crisis = CrisisSession(
            user_email=user_email,
            trigger_reason=payload.trigger_reason,
            breathing_completed=payload.breathing_completed,
            micro_tasks_generated=payload.micro_tasks_generated,
            duration_seconds=payload.duration_seconds,
            resolved=payload.resolved,
        )

        saved = save_crisis_session(crisis, user_email)
        if saved:
            return {"status": "ok", "session": saved.model_dump()}
        return {"status": "ok", "session": crisis.model_dump()}

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"❌ Error in crisis mode start: {e}")
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/crisis-mode/decompose")
async def decompose_tasks_for_crisis(
    authorization: str = Header(None),
):
    """Descompone las tareas pendientes del usuario en micro-tareas de 5 minutos.
    
    Usa las tareas pendientes reales del usuario y las fragmenta en pasos
    pequeños y manejables para reducir la ansiedad por overwhelm.
    """
    try:
        session = _require_session(authorization)
        user_email = session.email

        # Cargar tareas pendientes del usuario
        user_tasks = load_tasks(user_email)
        pending = [
            t for t in user_tasks
            if t.status in (TaskStatus.PENDING, TaskStatus.IN_PROGRESS)
        ]

        if not pending:
            return {
                "micro_tasks": [
                    {
                        "title": "🎉 ¡No tienes tareas pendientes!",
                        "estimated_minutes": 5,
                        "original_task_id": None,
                        "original_task_title": None,
                    }
                ],
                "message": "No hay tareas pendientes. ¡Buen trabajo!",
            }

        # Generar micro-tareas a partir de las pendientes (sin IA, rápido y offline-friendly)
        micro_tasks = []
        for task in pending[:5]:  # Máximo 5 tareas para no abrumar
            # Cada tarea se divide en micro-pasos de 5 min
            steps = _generate_micro_steps(task)
            micro_tasks.extend(steps)

        return {
            "micro_tasks": [mt.model_dump() for mt in micro_tasks[:8]],  # Máximo 8 micro-tareas
            "original_tasks_count": len(pending),
            "message": "Solo haz la primera. Nada más. 💪",
        }

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"❌ Error decomposing tasks: {e}")
        raise HTTPException(status_code=500, detail=str(e))


def _generate_micro_steps(task: Task) -> list:
    """Genera micro-pasos de 5 minutos a partir de una tarea.
    
    Heurística sencilla (sin llamada a IA) para ser instantáneo y offline.
    """
    title = task.title.strip()
    course = task.course.strip() if task.course else ""
    prefix = f"[{course}] " if course else ""

    steps = [
        MicroTask(
            title=f"{prefix}Abrir materiales de '{title}'",
            estimated_minutes=3,
            original_task_id=task.id,
            original_task_title=title,
        ),
        MicroTask(
            title=f"{prefix}Escribir 3 ideas clave sobre '{title}'",
            estimated_minutes=5,
            original_task_id=task.id,
            original_task_title=title,
        ),
    ]

    # Si la tarea tiene notas, añadir un paso de revisión
    if task.notes:
        steps.append(
            MicroTask(
                title=f"{prefix}Revisar tus notas anteriores de '{title}'",
                estimated_minutes=5,
                original_task_id=task.id,
                original_task_title=title,
            )
        )

    return steps


# ════════════════════════════════════════════════════════════════════
# ACCOUNT MANAGEMENT: Eliminación de cuenta (GDPR compliance)
#   DELETE /account → Elimina la cuenta del usuario y todos sus datos
# ════════════════════════════════════════════════════════════════════

@app.delete("/account", status_code=200)
def delete_account(authorization: Optional[str] = Header(None)):
    """Elimina permanentemente la cuenta del usuario y todos sus datos asociados.
    
    **GDPR COMPLIANCE**: Esta operación es irreversible y elimina:
    - Información del usuario
    - Todas las tareas
    - Todos los recordatorios
    - Entradas de calendario
    - Sesiones de enfoque
    - Estadísticas del usuario
    - Logros y retroalimentación
    - Suscripciones push
    - Todo historial personal
    
    **Headers requeridos:**
    - Authorization: Bearer <session_token>
    
    **Respuesta:** 200 OK si la eliminación fue exitosa
    
    **Errores:**
    - 401: Sesión inválida
    - 500: Error en la eliminación
    """
    try:
        # Validar sesión
        session = _require_session(authorization)
        user_email = session.email
        
        logger.warning(f"🗑️ INICANDO ELIMINACIÓN DE CUENTA PARA: {user_email}")
        
        # 1. Eliminar todas las tareas
        try:
            tasks = load_tasks(user_email)
            if tasks:
                save_tasks([], user_email=user_email)
                logger.info(f"  ✓ Eliminadas {len(tasks)} tareas")
        except Exception as e:
            logger.warning(f"  ⚠️ Error al eliminar tareas: {e}")
        
        # 2. Eliminar todos los recordatorios
        try:
            reminders = load_reminders(user_email)
            if reminders:
                save_reminders([], user_email=user_email)
                logger.info(f"  ✓ Eliminados {len(reminders)} recordatorios")
        except Exception as e:
            logger.warning(f"  ⚠️ Error al eliminar recordatorios: {e}")
        
        # 3. Eliminar entradas de calendario
        try:
            schedule = load_schedule(user_email)
            if schedule:
                save_schedule([], user_email=user_email)
                logger.info(f"  ✓ Eliminadas {len(schedule)} entradas de calendario")
        except Exception as e:
            logger.warning(f"  ⚠️ Error al eliminar calendario: {e}")
        
        # 4. Eliminar sesiones de enfoque
        try:
            from .supabase_client import get_supabase_client
            supabase = get_supabase_client()
            if supabase:
                supabase.table('focus_sessions').delete().eq('user_email', user_email).execute()
                logger.info(f"  ✓ Eliminadas sesiones de enfoque")
        except Exception as e:
            logger.warning(f"  ⚠️ Error al eliminar sesiones de enfoque: {e}")
        
        # 5. Eliminar estadísticas de usuario
        try:
            from .supabase_client import get_supabase_client
            supabase = get_supabase_client()
            if supabase:
                supabase.table('user_stats').delete().eq('user_email', user_email).execute()
                logger.info(f"  ✓ Eliminadas estadísticas de usuario")
        except Exception as e:
            logger.warning(f"  ⚠️ Error al eliminar estadísticas: {e}")
        
        # 6. Eliminar logros
        try:
            from .supabase_client import get_supabase_client
            supabase = get_supabase_client()
            if supabase:
                supabase.table('user_achievements').delete().eq('user_email', user_email).execute()
                logger.info(f"  ✓ Eliminados logros")
        except Exception as e:
            logger.warning(f"  ⚠️ Error al eliminar logros: {e}")
        
        # 7. Eliminar retroalimentación
        try:
            from .supabase_client import get_supabase_client
            supabase = get_supabase_client()
            if supabase:
                supabase.table('user_feedback').delete().eq('user_email', user_email).execute()
                logger.info(f"  ✓ Eliminada retroalimentación")
        except Exception as e:
            logger.warning(f"  ⚠️ Error al eliminar retroalimentación: {e}")
        
        # 8. Eliminar suscripciones push
        try:
            from .supabase_client import get_supabase_client
            supabase = get_supabase_client()
            if supabase:
                supabase.table('push_subscriptions').delete().eq('user_email', user_email).execute()
                logger.info(f"  ✓ Eliminadas suscripciones push")
        except Exception as e:
            logger.warning(f"  ⚠️ Error al eliminar suscripciones push: {e}")
        
        # 9. Eliminar notificaciones programadas
        try:
            from .supabase_client import get_supabase_client
            supabase = get_supabase_client()
            if supabase:
                supabase.table('scheduled_notifications').delete().eq('user_email', user_email).execute()
                logger.info(f"  ✓ Eliminadas notificaciones programadas")
        except Exception as e:
            logger.warning(f"  ⚠️ Error al eliminar notificaciones: {e}")
        
        # 10. Eliminar sesiones de crisis
        try:
            from .supabase_client import get_supabase_client
            supabase = get_supabase_client()
            if supabase:
                supabase.table('crisis_sessions').delete().eq('user_email', user_email).execute()
                logger.info(f"  ✓ Eliminadas sesiones de crisis")
        except Exception as e:
            logger.warning(f"  ⚠️ Error al eliminar sesiones de crisis: {e}")
        
        # 11. Eliminar entradas de energía
        try:
            from .supabase_client import get_supabase_client
            supabase = get_supabase_client()
            if supabase:
                supabase.table('energy_entries').delete().eq('user_email', user_email).execute()
                logger.info(f"  ✓ Eliminadas entradas de energía")
        except Exception as e:
            logger.warning(f"  ⚠️ Error al eliminar entradas de energía: {e}")
        
        # 12. Invalidar sesión
        try:
            invalidate_session(session.session_token)
            logger.info(f"  ✓ Sesión invalidada")
        except Exception as e:
            logger.warning(f"  ⚠️ Error al invalidar sesión: {e}")
        
        # 13. Eliminar usuario de la tabla users (último paso)
        try:
            from .supabase_client import get_supabase_client
            supabase = get_supabase_client()
            if supabase:
                supabase.table('users').delete().eq('email', user_email).execute()
                logger.info(f"  ✓ Eliminado usuario de tabla users")
        except Exception as e:
            logger.warning(f"  ⚠️ Error al eliminar usuario: {e}")
        
        logger.warning(f"✅ CUENTA ELIMINADA COMPLETAMENTE: {user_email}")
        
        return {
            "success": True,
            "message": "Tu cuenta ha sido eliminada permanentemente. Todos tus datos han sido borrados.",
            "email": user_email
        }
        
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"❌ Error deleting account: {e}", exc_info=True)
        raise HTTPException(status_code=500, detail="Error al eliminar la cuenta. Intenta nuevamente.")